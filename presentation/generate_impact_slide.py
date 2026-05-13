"""
Standalone generator — produces a single .pptx file containing ONLY the
Real-World Impact slide. Before/after consumption data from Karuna Hostel.
Copy-paste friendly into any deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
from PIL import Image

# ---------- PALETTE ----------
NAVY       = RGBColor(0x0A, 0x25, 0x40)
TEAL       = RGBColor(0x00, 0xC2, 0xB8)
TEAL_DARK  = RGBColor(0x00, 0x8C, 0x85)
AMBER      = RGBColor(0xF4, 0xA2, 0x61)
AMBER_DARK = RGBColor(0xD4, 0x84, 0x41)
ALERT      = RGBColor(0xE6, 0x39, 0x46)
SUCCESS    = RGBColor(0x2A, 0x9D, 0x8F)
BG_LIGHT   = RGBColor(0xF8, 0xF9, 0xFB)
DIVIDER    = RGBColor(0xE5, 0xE8, 0xEC)
TEXT_DARK  = RGBColor(0x2A, 0x2E, 0x37)
TEXT_MID   = RGBColor(0x5B, 0x64, 0x72)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

HERE = Path(__file__).parent
DEC_SHOT = HERE / "assets" / "screenshots" / "dec_consumption_cropped.png"
JAN_SHOT = HERE / "assets" / "screenshots" / "jan_consumption.png"


# ---------- HELPERS ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_rounded(slide, x, y, w, h, fill, line=None, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color=DIVIDER, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    if h:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    return slide.shapes.add_picture(str(path), x, y)


def section_pill(slide, x, y, w, text, *, fill, color=WHITE):
    """Small accent pill used as a column header."""
    add_rounded(slide, x, y, w, Inches(0.4), fill, radius=0.4)
    add_text(slide, x + Inches(0.2), y + Inches(0.02), w - Inches(0.4), Inches(0.36),
             text, size=11, bold=True, color=color,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ---------- DETECT ACTUAL IMAGE ASPECTS ----------
dec_img = Image.open(DEC_SHOT)
jan_img = Image.open(JAN_SHOT)
dec_aspect = dec_img.width / dec_img.height  # ~1.82
jan_aspect = jan_img.width / jan_img.height  # ~1.85


# ---------- SLIDE ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)

# --- Header ---
add_text(s, Inches(0.6), Inches(0.35), Inches(11), Inches(0.3),
         "FIELD RESULTS", size=11, bold=True, color=SUCCESS)
add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
         "Real-world impact at Karuna Hostel — G Block",
         size=28, bold=True, color=NAVY)
add_line(s, Inches(0.6), Inches(1.35), Inches(12.7), Inches(1.35), DIVIDER, 1.0)

# --- Subtitle ---
add_text(s, Inches(0.6), Inches(1.5), Inches(12.15), Inches(0.4),
         "Before and after device deployment — 19 December 2025",
         size=14, color=TEXT_MID, line_spacing=1.3)

# ============================================================
# TWO SIDE-BY-SIDE CHARTS
# ============================================================
# Equal image heights so they look balanced side-by-side.
img_h = Inches(3.35)
img_y = Inches(2.35)

# Use each image's natural aspect so nothing gets distorted
dec_w = Inches(3.35 * dec_aspect)     # ~6.09"
jan_w = Inches(3.35 * jan_aspect)     # ~6.19"

gap = Inches(0.3)
total_w = dec_w + gap + jan_w
left_margin = (SLIDE_W - total_w) / 2  # center the pair horizontally

dec_x = left_margin
jan_x = dec_x + dec_w + gap

# ---- Column header pills (above each image) ----
section_pill(s, dec_x, Inches(1.95), dec_w,
             "DECEMBER 2025  ·  DEVICE DEPLOYED 19 DEC", fill=AMBER_DARK)
section_pill(s, jan_x, Inches(1.95), jan_w,
             "JANUARY 2026  ·  FIRST FULL MONTH", fill=TEAL_DARK)

# ---- Dec screenshot (framed) ----
# shadow
add_rect(s, dec_x + Emu(55000), img_y + Emu(55000), dec_w, img_h,
         RGBColor(0xCF, 0xD5, 0xDC))
# frame
add_rounded(s, dec_x, img_y, dec_w, img_h, WHITE, line=DIVIDER, radius=0.015)
# image
pad = Inches(0.06)
add_image(s, DEC_SHOT, dec_x + pad, img_y + pad,
          w=dec_w - pad * 2, h=img_h - pad * 2)

# ---- Jan screenshot (framed) ----
add_rect(s, jan_x + Emu(55000), img_y + Emu(55000), jan_w, img_h,
         RGBColor(0xCF, 0xD5, 0xDC))
add_rounded(s, jan_x, img_y, jan_w, img_h, WHITE, line=DIVIDER, radius=0.015)
add_image(s, JAN_SHOT, jan_x + pad, img_y + pad,
          w=jan_w - pad * 2, h=img_h - pad * 2)

# ---- Captions under each chart ----
caption_y = img_y + img_h + Inches(0.2)

# Dec caption
add_text(s, dec_x, caption_y, dec_w, Inches(0.35),
         "Total: 974.39 kL  ·  transition month",
         size=13, bold=True, color=NAVY)
add_text(s, dec_x, caption_y + Inches(0.4), dec_w, Inches(0.8),
         "Offline data gaps and accumulated spikes before deployment. Clean consistent readings from 20 Dec onward — the new device in action.",
         size=11, color=TEXT_MID, line_spacing=1.4)

# Jan caption
add_text(s, jan_x, caption_y, jan_w, Inches(0.35),
         "Total: 887.55 kL  ·  first full month",
         size=13, bold=True, color=NAVY)
add_text(s, jan_x, caption_y + Inches(0.4), jan_w, Inches(0.8),
         "28 consecutive days of continuous data capture. No offline gaps, no accumulation spikes, no missing data.",
         size=11, color=TEXT_MID, line_spacing=1.4)

# ---- Bottom takeaway pill ----
takeaway_y = Inches(7.1)
add_rounded(s, Inches(0.6), takeaway_y, Inches(12.15), Inches(0.4),
            NAVY, radius=0.4)
add_text(s, Inches(0.9), takeaway_y, Inches(11.5), Inches(0.4),
         "One device, deployed once — complete site visibility for every day that followed.",
         size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

# --- Speaker notes ---
notes = s.notes_slide.notes_text_frame
notes.text = """Bring the room into the data. Speak slowly — this slide is evidence.

'These are real consumption readings from G Block, Karuna Hostel — a site where we deployed the new device on 19 December 2025.

On the left — December 2025. Look at the pattern. Days 1 through 19 show a mix of grey and teal bars. Grey means offline data — readings that couldn't reach the cloud when they happened. The big spike on day 6 is an accumulation — data that piled up during an outage. From day 15 through 19, there's almost nothing visible — those are days where the old device lost data entirely. Then something changes on 20 December. Clean, consistent teal bars from that point onward. That's the new device in action.

On the right — January 2026. The first full month with the new device. Twenty-eight consecutive days of continuous data capture. No offline gaps. No accumulation spikes. No missing data.

One device, deployed once, on one day. Complete site visibility every day that followed.'

Timing: 75 seconds. Point to the transition around day 19-20 on the Dec chart — let the audience see it."""


# ---------- SAVE ----------
OUT = Path(__file__).parent / "Real_World_Impact.pptx"
prs.save(OUT)
print(f"Generated: {OUT}")
print(f"Size: {OUT.stat().st_size // 1024} KB")
print(f"Dec aspect: {dec_aspect:.3f}:1, Jan aspect: {jan_aspect:.3f}:1 (both preserved)")
