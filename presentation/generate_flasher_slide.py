"""
Standalone generator — produces a single .pptx file containing ONLY the
Firmware Flasher slide.
Layout: screenshot on the LEFT with natural aspect ratio preserved,
        benefits stacked on the RIGHT.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
from PIL import Image

# ---------- PALETTE ----------
NAVY       = RGBColor(0x0A, 0x25, 0x40)
TEAL       = RGBColor(0x00, 0xC2, 0xB8)
TEAL_DARK  = RGBColor(0x00, 0x8C, 0x85)
AMBER      = RGBColor(0xF4, 0xA2, 0x61)
AMBER_DARK = RGBColor(0xD4, 0x84, 0x41)
BG_LIGHT   = RGBColor(0xF8, 0xF9, 0xFB)
DIVIDER    = RGBColor(0xE5, 0xE8, 0xEC)
TEXT_DARK  = RGBColor(0x2A, 0x2E, 0x37)
TEXT_MID   = RGBColor(0x5B, 0x64, 0x72)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

HERE = Path(__file__).parent
SHOT_PATH = HERE / "assets" / "screenshots" / "flasher_tool.png"


# ---------- HELPERS ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_rounded(slide, x, y, w, h, fill, line=None, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color=DIVIDER, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    if h:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    return slide.shapes.add_picture(str(path), x, y)


# ---------- DETECT REAL IMAGE ASPECT ----------
img = Image.open(SHOT_PATH)
img_aspect = img.width / img.height  # ~2.28 for this screenshot

# ---------- SLIDE ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)

# --- Header ---
add_text(s, Inches(0.6), Inches(0.35), Inches(11), Inches(0.3),
         "FIRMWARE FLASHER", size=11, bold=True, color=TEAL_DARK)
add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
         "Initial firmware setup — no software required",
         size=30, bold=True, color=NAVY)
add_line(s, Inches(0.6), Inches(1.35), Inches(12.7), Inches(1.35), DIVIDER, 1.0)

# --- Subtitle / tagline ---
add_text(s, Inches(0.6), Inches(1.55), Inches(12.15), Inches(0.4),
         "A browser-based flasher — any user can install firmware on an ESP32 with a single click.",
         size=13, color=TEXT_MID, line_spacing=1.3)

# ============================================================
# MAIN CONTENT AREA  —  2.15" to 6.75"  (4.6" tall)
# ============================================================
content_top = Inches(2.15)
content_bottom = Inches(6.75)
content_h = content_bottom - content_top

# ---- LEFT: SCREENSHOT (correct aspect) ----
shot_w = Inches(7.7)                        # 7.7" wide
shot_h = Inches(7.7 / img_aspect)           # preserves 2.28:1 natural ratio
# Vertically centered in the content area
shot_x = Inches(0.6)
shot_y = content_top + (content_h - shot_h) / 2

# Outer soft shadow
add_rect(s, shot_x + Emu(55000), shot_y + Emu(55000),
         shot_w, shot_h, RGBColor(0xCF, 0xD5, 0xDC))

# White frame/padding so the dark screenshot has breathing room
frame_pad = Inches(0.08)
add_rounded(s, shot_x, shot_y, shot_w, shot_h, WHITE, line=DIVIDER, radius=0.015)

# Actual screenshot — natural aspect preserved
add_image(s, SHOT_PATH,
          shot_x + frame_pad, shot_y + frame_pad,
          w=shot_w - frame_pad * 2,
          h=shot_h - frame_pad * 2)

# ---- RIGHT: BENEFIT CARDS (stacked) ----
right_col_x = shot_x + shot_w + Inches(0.3)         # 8.6"
right_col_w = SLIDE_W - right_col_x - Inches(0.6)   # remainder, respecting margin

benefits = [
    ("01", "No software required",
     "Runs directly in Chrome or Edge using Web Serial — nothing to install.",
     TEAL),
    ("02", "One-click install",
     "Bootloader, partition table, and firmware — flashed together in one operation.",
     TEAL_DARK),
    ("03", "Anyone can deploy",
     "No Arduino IDE, no command line, no developer required.",
     AMBER),
]

# 3 cards vertically stacked, total height fits content area
card_gap = Inches(0.15)
card_count = len(benefits)
card_h = (content_h - card_gap * (card_count - 1)) / card_count

for i, (num, title, body, accent) in enumerate(benefits):
    cy = content_top + i * (card_h + card_gap)
    cx = right_col_x
    cw = right_col_w
    # shadow
    add_rect(s, cx + Emu(40000), cy + Emu(40000), cw, card_h,
             RGBColor(0xDD, 0xE1, 0xE7))
    # body
    add_rect(s, cx, cy, cw, card_h, WHITE, line=DIVIDER)
    # left accent bar
    add_rect(s, cx, cy, Inches(0.12), card_h, accent)
    # Number (small, colored)
    add_text(s, cx + Inches(0.35), cy + Inches(0.2), Inches(0.5), Inches(0.3),
             num, size=11, bold=True, color=accent)
    # Title
    add_text(s, cx + Inches(0.35), cy + Inches(0.45), cw - Inches(0.5), Inches(0.5),
             title, size=16, bold=True, color=NAVY, line_spacing=1.15)
    # Body
    add_text(s, cx + Inches(0.35), cy + Inches(0.9), cw - Inches(0.5),
             card_h - Inches(1.0),
             body, size=11.5, color=TEXT_MID, line_spacing=1.4)

# --- Bottom takeaway pill ---
takeaway_y = Inches(6.95)
add_rounded(s, Inches(0.6), takeaway_y, Inches(12.15), Inches(0.45),
            NAVY, radius=0.4)
add_text(s, Inches(0.9), takeaway_y, Inches(11.5), Inches(0.45),
         "Any user can flash an ESP32 in under a minute — from a browser, with zero setup.",
         size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

# --- Speaker notes ---
notes = s.notes_slide.notes_text_frame
notes.text = """'We also built a web-based firmware flasher to make the initial setup of every device as easy as possible.

Here's how it works. Any user opens this page in Chrome or Edge. They plug in the ESP32 via USB. They click Install Firmware. The browser installs the full firmware — bootloader, partition table, and application — in one shot, directly through the Web Serial API. No Arduino IDE, no command line, no local software of any kind.

Three reasons this matters.

One — no software required. It runs entirely in the browser. There is nothing to install.

Two — it's a one-click install. The full firmware goes on in a single operation.

Three — anyone can deploy a device. You don't need to be a developer. Operations, customer, or new hire — anyone with a browser and a USB cable.

Any user can flash an ESP32 in under a minute, from a browser, with zero setup.'

Timing: 60 seconds."""


# ---------- SAVE ----------
OUT = Path(__file__).parent / "Firmware_Flasher.pptx"
prs.save(OUT)
print(f"Generated: {OUT}")
print(f"Size: {OUT.stat().st_size // 1024} KB")
print(f"Screenshot aspect: {img_aspect:.3f}:1 (natural ratio preserved)")
