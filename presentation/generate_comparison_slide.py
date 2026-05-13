"""
Standalone generator — produces a single .pptx file containing ONLY the
Aquagen vs Teltonika comparison slide.
Copy-paste-friendly into any other presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ---------- PALETTE ----------
NAVY       = RGBColor(0x0A, 0x25, 0x40)
TEAL       = RGBColor(0x00, 0xC2, 0xB8)
TEAL_DARK  = RGBColor(0x00, 0x8C, 0x85)
AMBER_DARK = RGBColor(0xD4, 0x84, 0x41)
BG_LIGHT   = RGBColor(0xF8, 0xF9, 0xFB)
DIVIDER    = RGBColor(0xE5, 0xE8, 0xEC)
TEXT_MID   = RGBColor(0x5B, 0x64, 0x72)
TEXT_LIGHT = RGBColor(0x95, 0x9C, 0xA8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------- HELPERS ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_rounded(slide, x, y, w, h, fill, line=None, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color=DIVIDER, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


# ---------- SLIDE ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)

# Page header
add_text(s, Inches(0.6), Inches(0.35), Inches(11), Inches(0.3),
         "IN-HOUSE vs EXTERNAL", size=11, bold=True, color=TEAL)
add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
         "Aquagen vs Teltonika — feature-by-feature comparison",
         size=30, bold=True, color=NAVY)
add_line(s, Inches(0.6), Inches(1.35), Inches(12.7), Inches(1.35), DIVIDER, 1.0)

# Subtitle
add_text(s, Inches(0.6), Inches(1.55), Inches(12.15), Inches(0.4),
         "Feature parity on core capabilities — with ownership, direct Azure integration, and no vendor subscription.",
         size=12, color=TEXT_MID, line_spacing=1.3)

# Two product cards
card_w = Inches(6.05)
card_h = Inches(4.95)
card_y = Inches(2.1)
card_gap = Inches(0.15)
left_card_x = Inches(0.6)
right_card_x = left_card_x + card_w + card_gap


def product_card(slide, x, y, w, h, *, title_text, subtitle_text, tag_text,
                 accent_color, row_items):
    # shadow
    add_rect(slide, x + Emu(45000), y + Emu(45000), w, h, RGBColor(0xDD, 0xE1, 0xE7))
    # body
    add_rect(slide, x, y, w, h, WHITE, line=DIVIDER)
    # header
    add_rect(slide, x, y, w, Inches(1.0), accent_color)
    add_text(slide, x + Inches(0.4), y + Inches(0.15), w - Inches(0.8), Inches(0.35),
             tag_text.upper(), size=10, bold=True, color=WHITE)
    add_text(slide, x + Inches(0.4), y + Inches(0.45), w - Inches(0.8), Inches(0.5),
             title_text, size=22, bold=True, color=WHITE)
    add_text(slide, x + Inches(0.4), y + Inches(0.78), w - Inches(0.8), Inches(0.3),
             subtitle_text, size=10, color=RGBColor(0xE4, 0xF8, 0xF7))

    row_y = y + Inches(1.2)
    row_h = Inches(0.38)
    for i, (label, value) in enumerate(row_items):
        if i % 2 == 0:
            add_rect(slide, x + Inches(0.2), row_y, w - Inches(0.4), row_h,
                     RGBColor(0xFA, 0xFB, 0xFC))
        add_text(slide, x + Inches(0.35), row_y + Inches(0.09),
                 Inches(2.1), Inches(0.25),
                 label, size=9.5, bold=True, color=TEXT_MID)
        add_text(slide, x + Inches(2.55), row_y + Inches(0.07),
                 w - Inches(2.75), Inches(0.3),
                 value, size=10.5, color=NAVY, line_spacing=1.25)
        row_y += row_h


aquagen_rows = [
    ("Platform",          "ESP32-WROOM-32 · Custom hardware"),
    ("Connectivity",      "WiFi + 4G LTE (A7670C cellular)"),
    ("Industrial bus",    "RS485 Modbus RTU · 10 sensors per unit"),
    ("Offline storage",   "SD Card · automatic + chronological replay"),
    ("Timestamps",        "DS3231 RTC · accurate even offline"),
    ("Cloud integration", "Azure IoT Hub · MQTT/TLS direct"),
    ("Remote management", "Azure Device Twin · built-in · no subscription"),
    ("Configuration",     "Custom web portal · branded UI"),
    ("OTA updates",       "Signed firmware · dual-slot rollback"),
    ("Customization",     "Full source control · any change possible"),
    ("Vendor dependency", "None · built in-house at FluxGen"),
]
product_card(s, left_card_x, card_y, card_w, card_h,
             title_text="Aquagen",
             subtitle_text="FluxGen in-house IoT gateway",
             tag_text="IN-HOUSE · BUILT BY US",
             accent_color=TEAL_DARK, row_items=aquagen_rows)

teltonika_rows = [
    ("Platform",          "RUT956 / TRB246 · vendor-built router"),
    ("Connectivity",      "WiFi + 4G LTE · Dual-SIM failover"),
    ("Industrial bus",    "RS232 + RS485 · digital I/O + relay"),
    ("Offline storage",   "microSD (requires configuration)"),
    ("Timestamps",        "System clock · NTP-synced"),
    ("Cloud integration", "Via MQTT broker or Teltonika RMS"),
    ("Remote management", "Teltonika RMS · paid subscription"),
    ("Configuration",     "RutOS web UI · vendor interface"),
    ("OTA updates",       "Vendor-supplied firmware updates"),
    ("Customization",     "Limited to RutOS parameters"),
    ("Vendor dependency", "Teltonika (Lithuania-based vendor)"),
]
product_card(s, right_card_x, card_y, card_w, card_h,
             title_text="Teltonika",
             subtitle_text="External procurement · RUT956 / TRB246 series",
             tag_text="EXTERNAL · WE BUY",
             accent_color=AMBER_DARK, row_items=teltonika_rows)

# Takeaway pill
takeaway_y = Inches(7.15)
add_rounded(s, Inches(0.6), takeaway_y - Inches(0.03), Inches(12.15), Inches(0.4),
            NAVY, radius=0.4)
add_text(s, Inches(0.9), takeaway_y, Inches(11.5), Inches(0.35),
         "We match them on core features — and own the stack, the roadmap, and the cloud integration.",
         size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

# Speaker notes
notes = s.notes_slide.notes_text_frame
notes.text = """Let me put this in context. Teltonika is what we currently buy externally — specifically their RUT956 and TRB246 series. Good products. Widely used.

On the left — Aquagen, our in-house build. On the right — Teltonika, what we buy.

Look at the feature rows. Connectivity — both do WiFi plus 4G LTE. Industrial bus — both do RS485 Modbus. Offline storage — both support SD-card caching. Cloud — Teltonika integrates via their RMS subscription or an MQTT broker; Aquagen connects directly to Azure IoT Hub.

The real differences are at the bottom. Customization — Aquagen is fully in-house; we can change anything. Teltonika is locked to RutOS parameters. Vendor dependency — Aquagen has none; Teltonika ties us to a Lithuania-based vendor and their subscription model.

Bottom line — we match them on core features, but we own the stack, the roadmap, and the cloud integration."""


# ---------- SAVE ----------
OUT = Path(__file__).parent / "Aquagen_vs_Teltonika_Comparison.pptx"
prs.save(OUT)
print(f"Generated: {OUT}")
print(f"Size: {OUT.stat().st_size // 1024} KB")
