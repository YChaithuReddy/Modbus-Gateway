"""
FluxGen IoT Gateway — v2.0 Pitch Deck Generator
Executive 10-slide deck for 10-minute management presentation.
Style: McKinsey/Apple-inspired, navy + teal + amber.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree

# ---------- DESIGN SYSTEM ----------
NAVY       = RGBColor(0x0A, 0x25, 0x40)
NAVY_DARK  = RGBColor(0x06, 0x17, 0x2A)
TEAL       = RGBColor(0x00, 0xC2, 0xB8)
TEAL_DARK  = RGBColor(0x00, 0x8C, 0x85)
AMBER      = RGBColor(0xF4, 0xA2, 0x61)
AMBER_DARK = RGBColor(0xD4, 0x84, 0x41)
BG_LIGHT   = RGBColor(0xF8, 0xF9, 0xFB)
DIVIDER    = RGBColor(0xE5, 0xE8, 0xEC)
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MID   = RGBColor(0x5B, 0x64, 0x72)
TEXT_LIGHT = RGBColor(0x95, 0x9C, 0xA8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS    = RGBColor(0x2A, 0x9D, 0x8F)
ALERT      = RGBColor(0xE6, 0x39, 0x46)

FONT = "Calibri"
FONT_HEAD = "Calibri"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# ---------- HELPERS ----------
def add_rect(slide, x, y, w, h, fill, line=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    if not shadow:
        shape.shadow.inherit = False
    return shape


def add_rounded(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_DARK,
             font=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT_DARK,
                bullet_color=None, font=FONT, line_spacing=1.35):
    """items: list of (bold_prefix, rest) or strings. Custom dot bullets."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # bullet dot
        r0 = p.add_run()
        r0.text = "●  "
        r0.font.name = font
        r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color or color
        r0.font.bold = True
        if isinstance(item, tuple):
            bold_part, rest = item
            r1 = p.add_run(); r1.text = bold_part
            r1.font.name = font; r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = color
            r2 = p.add_run(); r2.text = rest
            r2.font.name = font; r2.font.size = Pt(size); r2.font.color.rgb = color
        else:
            r1 = p.add_run(); r1.text = item
            r1.font.name = font; r1.font.size = Pt(size); r1.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color=DIVIDER, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def page_header(slide, kicker, title, *, kicker_color=TEAL):
    """Top band: teal kicker + bold navy title + divider."""
    add_text(slide, Inches(0.6), Inches(0.35), Inches(11), Inches(0.3),
             kicker.upper(), size=11, bold=True, color=kicker_color)
    add_text(slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    add_line(slide, Inches(0.6), Inches(1.35), Inches(12.7), Inches(1.35), DIVIDER, 1.0)


def page_footer(slide, text_left="FluxGen · IoT Platform", text_right=""):
    add_text(slide, Inches(0.6), Inches(7.1), Inches(6), Inches(0.3),
             text_left, size=9, color=TEXT_LIGHT)
    if text_right:
        add_text(slide, Inches(7.3), Inches(7.1), Inches(5.4), Inches(0.3),
                 text_right, size=9, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def card(slide, x, y, w, h, *, accent=TEAL, headline="", body_lines=None, icon=None):
    """Card with accent top band, big headline, body bullets."""
    # shadow
    add_rect(slide, x + Emu(36000), y + Emu(36000), w, h, RGBColor(0xE0, 0xE3, 0xE8))
    # card body
    box = add_rect(slide, x, y, w, h, WHITE, line=DIVIDER)
    # accent top
    add_rect(slide, x, y, w, Inches(0.12), accent)
    inner_x = x + Inches(0.3)
    inner_w = w - Inches(0.6)
    cursor_y = y + Inches(0.3)
    if icon:
        add_text(slide, inner_x, cursor_y, inner_w, Inches(0.6),
                 icon, size=32, color=accent, bold=True)
        cursor_y += Inches(0.65)
    add_text(slide, inner_x, cursor_y, inner_w, Inches(0.5),
             headline, size=16, bold=True, color=NAVY, line_spacing=1.2)
    cursor_y += Inches(0.55)
    if body_lines:
        add_text(slide, inner_x, cursor_y, inner_w, h - (cursor_y - y) - Inches(0.2),
                 "\n".join(body_lines), size=12, color=TEXT_MID, line_spacing=1.35)


def pill(slide, x, y, w, h, text, *, fill=TEAL, color=WHITE, size=10):
    pill_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill_shape.adjustments[0] = 0.5
    pill_shape.fill.solid()
    pill_shape.fill.fore_color.rgb = fill
    pill_shape.line.fill.background()
    pill_shape.shadow.inherit = False
    tf = pill_shape.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = color


# =========================================================
#  SLIDE 1 — TITLE
# =========================================================
s1 = prs.slides.add_slide(BLANK)
# Left navy panel
add_rect(s1, Inches(0), Inches(0), Inches(5.5), SLIDE_H, NAVY)
# Teal accent bar
add_rect(s1, Inches(5.5), Inches(0), Inches(0.08), SLIDE_H, TEAL)
# Right light panel
add_rect(s1, Inches(5.58), Inches(0), Inches(7.76), SLIDE_H, BG_LIGHT)

# Left side content
add_text(s1, Inches(0.7), Inches(0.7), Inches(4.5), Inches(0.4),
         "FLUXGEN · INTERNAL", size=11, bold=True, color=TEAL)
add_text(s1, Inches(0.7), Inches(2.4), Inches(4.7), Inches(1.2),
         "FluxGen", size=54, bold=True, color=WHITE, line_spacing=1.0)
add_text(s1, Inches(0.7), Inches(3.3), Inches(4.7), Inches(1.2),
         "IoT Gateway", size=54, bold=True, color=WHITE, line_spacing=1.0)
add_text(s1, Inches(0.7), Inches(4.5), Inches(4.7), Inches(0.5),
         "v2.0 · Proposal", size=20, color=TEAL)

add_rect(s1, Inches(0.7), Inches(5.25), Inches(1.2), Inches(0.04), TEAL)
add_text(s1, Inches(0.7), Inches(5.4), Inches(4.7), Inches(0.4),
         "Y CHAITHANYA REDDY", size=12, bold=True, color=WHITE)
add_text(s1, Inches(0.7), Inches(5.75), Inches(4.7), Inches(0.35),
         "IoT Platform — Firmware & Cloud", size=11, color=TEXT_LIGHT)
add_text(s1, Inches(0.7), Inches(6.7), Inches(4.7), Inches(0.3),
         "April 2026", size=10, color=TEXT_LIGHT)

# Right side — tagline + bullets
add_text(s1, Inches(6.1), Inches(1.4), Inches(6.8), Inches(0.4),
         "THE PITCH IN ONE LINE", size=11, bold=True, color=AMBER_DARK)
add_text(s1, Inches(6.1), Inches(1.8), Inches(6.8), Inches(2.4),
         "From field reliability to\nfleet intelligence.",
         size=34, bold=True, color=NAVY, line_spacing=1.1)

add_line(s1, Inches(6.1), Inches(4.2), Inches(12.7), Inches(4.2), DIVIDER, 0.75)

add_text(s1, Inches(6.1), Inches(4.4), Inches(6.8), Inches(0.4),
         "WHAT THIS DECK COVERS", size=10, bold=True, color=TEXT_MID)
add_bullets(s1, Inches(6.1), Inches(4.75), Inches(6.8), Inches(2.2), [
    ("The problem we solved — ", "field data loss across unreliable networks."),
    ("What v1.x delivers today — ", "a self-healing, offline-first edge gateway."),
    ("Why v2.0 matters — ", "scale from devices to platform."),
    ("The ask — ", "greenlight v2.0 and expanded scope."),
], size=13, color=TEXT_DARK, bullet_color=TEAL, line_spacing=1.5)

page_footer(s1, "FluxGen · Confidential · Internal Distribution", "Slide 1 of 10")

add_notes(s1, """Open strong — lock eye contact with CEO before first word.

'Good morning. In the next ten minutes I'm going to walk you through the IoT Gateway programme — the problem it solved, the business value v1.x is delivering today, and the case for v2.0.

By the end of this meeting I'll be asking for one decision: to greenlight v2.0 and expand the scope of the IoT Platform team. I'll make sure you have everything you need to make that call confidently.'

Pause 2 seconds. Click to next slide.""")


# =========================================================
#  SLIDE 2 — EXECUTIVE SUMMARY
# =========================================================
s2 = prs.slides.add_slide(BLANK)
page_header(s2, "Executive Summary", "The three things to remember")

# 3 cards
card_w = Inches(3.95); card_h = Inches(4.6); gap = Inches(0.2)
start_x = Inches(0.6); y = Inches(1.8)

# Card 1
card(s2, start_x, y, card_w, card_h, accent=TEAL, icon="01",
     headline="We solved field\ndata loss.",
     body_lines=[
         "Dual connectivity (WiFi + Cellular),",
         "SD-card offline caching, and auto",
         "self-healing — the gateway keeps",
         "data safe through network outages.",
         "",
         "Zero messages lost when the network",
         "drops. Proven in the field.",
     ])

# Card 2
card(s2, start_x + card_w + gap, y, card_w, card_h, accent=AMBER, icon="02",
     headline="We eliminated\ntruck rolls.",
     body_lines=[
         "Remote OTA via Azure IoT Hub.",
         "Remote sensor configuration.",
         "Remote reboot & diagnostics.",
         "",
         "Field service visits are no longer",
         "required for firmware, configuration,",
         "or recovery — at a fraction of",
         "imported PLC cost.",
     ])

# Card 3
card(s2, start_x + 2*(card_w + gap), y, card_w, card_h, accent=NAVY, icon="03",
     headline="We are ready\nto scale.",
     body_lines=[
         "v1.x is production-grade. The next",
         "lever is fleet intelligence — a",
         "centralised console, edge analytics,",
         "and security hardening for enterprise",
         "and government tenders.",
         "",
         "This is the v2.0 proposal.",
     ])

# Bottom callout
add_rounded(s2, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.48), NAVY)
add_text(s2, Inches(0.9), Inches(6.58), Inches(11.8), Inches(0.42),
         "The ask  →  Greenlight v2.0 · 16-week roadmap · 2 engineers · IoT Platform Tech Lead mandate",
         size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
page_footer(s2, "FluxGen · IoT Platform", "Slide 2 of 10")

add_notes(s2, """This is the TL;DR slide — if management only remembers three things, it's these.

'Three headlines before we go deep:

ONE — We solved field data loss. Our gateway keeps data safe through network outages that would have cost us messages before.

TWO — We eliminated truck rolls. Firmware, configuration and recovery all happen remotely, at a fraction of the cost of imported PLCs.

THREE — We are ready to scale. v1.x is production-grade. The next lever is not more device features — it is fleet intelligence.

The bar at the bottom summarises my ask. I'll unpack it on the final slide.'

Timing: 60 seconds. Do not linger.""")


# =========================================================
#  SLIDE 3 — THE PROBLEM
# =========================================================
s3 = prs.slides.add_slide(BLANK)
page_header(s3, "The Problem", "What was breaking before v1.x")

# Left — pain panel
left_x = Inches(0.6); left_w = Inches(6.0)
panel = add_rounded(s3, left_x, Inches(1.75), left_w, Inches(5.1), RGBColor(0xFD, 0xF3, 0xF3))
add_text(s3, left_x + Inches(0.4), Inches(1.95), left_w - Inches(0.8), Inches(0.35),
         "FIELD REALITY", size=11, bold=True, color=ALERT)
add_text(s3, left_x + Inches(0.4), Inches(2.3), left_w - Inches(0.8), Inches(0.6),
         "Data never made it back.", size=22, bold=True, color=NAVY, line_spacing=1.1)

add_bullets(s3, left_x + Inches(0.4), Inches(3.1), left_w - Inches(0.8), Inches(3.7), [
    ("Network drops → ", "telemetry lost forever. No replay, no recovery."),
    ("Every firmware update → ", "a truck roll to each plant site."),
    ("Sensor changes → ", "onsite laptop and a trained engineer."),
    ("No remote visibility → ", "customer calls before we know there is an issue."),
    ("Single connectivity path → ", "one failure mode stops reporting."),
], size=13, color=TEXT_DARK, bullet_color=ALERT, line_spacing=1.6)

# Right — business impact
right_x = Inches(6.9); right_w = Inches(5.85)
add_text(s3, right_x, Inches(1.85), right_w, Inches(0.4),
         "BUSINESS IMPACT", size=11, bold=True, color=AMBER_DARK)
add_text(s3, right_x, Inches(2.2), right_w, Inches(0.6),
         "The cost of not knowing.", size=22, bold=True, color=NAVY, line_spacing=1.1)

# 4 impact tiles 2x2
tile_w = Inches(2.85); tile_h = Inches(1.85); tile_gap = Inches(0.15)
impacts = [
    ("Compliance gaps", "CPCB/PCB reporting windows missed when data drops.", ALERT),
    ("Truck-roll cost", "Every firmware patch = engineer travel to every site.", AMBER_DARK),
    ("Trust erosion", "Customers discover outages before we do.", ALERT),
    ("Slow expansion", "Commissioning new sites required specialist onsite.", AMBER_DARK),
]
for i, (title, body, col) in enumerate(impacts):
    col_idx = i % 2; row_idx = i // 2
    tx = right_x + col_idx * (tile_w + tile_gap)
    ty = Inches(3.0) + row_idx * (tile_h + tile_gap)
    add_rect(s3, tx, ty, tile_w, tile_h, WHITE, line=DIVIDER)
    add_rect(s3, tx, ty, Inches(0.08), tile_h, col)
    add_text(s3, tx + Inches(0.22), ty + Inches(0.18), tile_w - Inches(0.3), Inches(0.4),
             title, size=13, bold=True, color=NAVY)
    add_text(s3, tx + Inches(0.22), ty + Inches(0.6), tile_w - Inches(0.3), tile_h - Inches(0.7),
             body, size=11, color=TEXT_MID, line_spacing=1.35)

page_footer(s3, "FluxGen · IoT Platform · The Problem", "Slide 3 of 10")

add_notes(s3, """Slow down here — this is the 'why should management care' slide.

'Before v1.x, the story in the field was simple: data did not always make it back.

When a cellular network dropped, telemetry from that window was lost permanently. There was no local buffer, no replay. For a customer submitting data to the Pollution Control Board, a silent outage is a compliance gap.

Every firmware update required a truck roll. Every sensor change required an engineer with a laptop onsite. And because we had no remote view, our customers often noticed a problem before we did.

The business impact is on the right — compliance risk, operational cost, trust erosion, and slow expansion. This is the mountain v1.x had to move.'

Timing: 75 seconds. Point at the right panel with flat hand — don't pick at bullets.""")


# =========================================================
#  SLIDE 4 — ROOT CAUSE
# =========================================================
s4 = prs.slides.add_slide(BLANK)
page_header(s4, "Root Cause", "Why traditional gateways fail in Indian field conditions")

# Intro line
add_text(s4, Inches(0.6), Inches(1.55), Inches(12), Inches(0.4),
         "Four architectural gaps — each one alone is enough to break reliability.",
         size=13, color=TEXT_MID)

# 2x2 diagnosis grid
grid_x = Inches(0.6); grid_y = Inches(2.15)
box_w = Inches(6.0); box_h = Inches(2.3); box_gap = Inches(0.15)
diag = [
    ("Single connectivity path",
     "Most gateways ship with WiFi OR cellular. A single failure mode stops reporting. No fallback when the primary link drops.",
     "No redundancy"),
    ("No local persistence",
     "When the network is down, telemetry is held in RAM for seconds, then overwritten. Multi-hour outages = permanent data loss.",
     "No buffer"),
    ("Physical-only configuration",
     "Sensor Modbus maps, addresses, and byte order are set with a laptop onsite. Every change needs an engineer in the field.",
     "No remote ops"),
    ("Black-box diagnostics",
     "No reliable way to see if a remote unit is healthy, let alone what it is doing. Issues surface only when the customer calls.",
     "No visibility"),
]
for i, (title, body, tag) in enumerate(diag):
    col_idx = i % 2; row_idx = i // 2
    bx = grid_x + col_idx * (box_w + box_gap)
    by = grid_y + row_idx * (box_h + box_gap)
    add_rect(s4, bx, by, box_w, box_h, WHITE, line=DIVIDER)
    # number badge
    badge = add_rect(s4, bx + Inches(0.3), by + Inches(0.3), Inches(0.5), Inches(0.5), NAVY)
    add_text(s4, bx + Inches(0.3), by + Inches(0.3), Inches(0.5), Inches(0.5),
             str(i + 1), size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title
    add_text(s4, bx + Inches(1.0), by + Inches(0.3), box_w - Inches(3.2), Inches(0.5),
             title, size=15, bold=True, color=NAVY)
    # pill tag
    pill(s4, bx + box_w - Inches(2.0), by + Inches(0.35), Inches(1.7), Inches(0.35),
         tag.upper(), fill=AMBER, color=NAVY_DARK, size=9)
    # body
    add_text(s4, bx + Inches(1.0), by + Inches(0.95), box_w - Inches(1.3), box_h - Inches(1.15),
             body, size=12, color=TEXT_MID, line_spacing=1.4)

page_footer(s4, "FluxGen · Root Cause Analysis", "Slide 4 of 10")

add_notes(s4, """This is the engineer-credibility slide. Speak crisply.

'Why do traditional gateways fail in the field? Four architectural gaps — and any one of them alone breaks reliability.

ONE, single connectivity path — WiFi or cellular, never both. A single failure mode stops the device.

TWO, no local persistence. Telemetry sits in RAM for seconds before it is overwritten. A multi-hour outage means permanent data loss.

THREE, physical-only configuration. Every sensor change requires a laptop onsite.

FOUR, black-box diagnostics. You learn about problems when a customer calls.

Our job in v1.x was to close all four — not just one.'

Timing: 60 seconds. If CTO asks why these are the four, say 'these are the ones we saw break in our own customer deployments'.""")


# =========================================================
#  SLIDE 5 — v1.x SOLUTIONS
# =========================================================
s5 = prs.slides.add_slide(BLANK)
page_header(s5, "What We Built — v1.x", "A self-healing edge gateway, production-grade today")

# Top row — 3 large cards
row1_y = Inches(1.7)
sol_w = Inches(4.0); sol_h = Inches(2.55); sol_gap = Inches(0.1)
solutions_row1 = [
    ("📡", "Dual Connectivity",
     "WiFi + Cellular (A7670C). Auto-failover between links.",
     "Works wherever the customer needs us."),
    ("💾", "SD-Card Offline Cache",
     "FAT32 on SD. Chronological replay. Power-loss safe.",
     "Zero data loss through outages."),
    ("🔄", "Auto Self-Healing",
     "Modem power-cycle on MQTT drop. 5-min cooldown logic.",
     "Recovers without a truck roll."),
]
for i, (icon, title, body, kicker) in enumerate(solutions_row1):
    sx = Inches(0.6) + i * (sol_w + sol_gap)
    add_rect(s5, sx + Emu(36000), row1_y + Emu(36000), sol_w, sol_h, RGBColor(0xE0, 0xE3, 0xE8))
    add_rect(s5, sx, row1_y, sol_w, sol_h, WHITE, line=DIVIDER)
    add_rect(s5, sx, row1_y, sol_w, Inches(0.1), TEAL)
    add_text(s5, sx + Inches(0.3), row1_y + Inches(0.25), sol_w - Inches(0.6), Inches(0.7),
             icon, size=28, color=TEAL)
    add_text(s5, sx + Inches(0.3), row1_y + Inches(0.95), sol_w - Inches(0.6), Inches(0.5),
             title, size=16, bold=True, color=NAVY)
    add_text(s5, sx + Inches(0.3), row1_y + Inches(1.4), sol_w - Inches(0.6), Inches(0.8),
             body, size=11, color=TEXT_MID, line_spacing=1.35)
    # kicker at bottom
    add_rect(s5, sx, row1_y + sol_h - Inches(0.45), sol_w, Inches(0.45), BG_LIGHT)
    add_text(s5, sx + Inches(0.3), row1_y + sol_h - Inches(0.4), sol_w - Inches(0.6), Inches(0.35),
             kicker, size=11, bold=True, color=TEAL_DARK, anchor=MSO_ANCHOR.MIDDLE)

# Bottom row — 2 cards (wider)
row2_y = Inches(4.45)
sol2_w = Inches(6.15); sol2_h = Inches(2.3); sol2_gap = Inches(0.15)
solutions_row2 = [
    ("☁️", "Remote OTA + Azure Device Twin",
     "Cloud-pushed firmware, config, and control. Signed boot.",
     "Fleet-wide updates, zero field visits."),
    ("🌐", "Web-Based Commissioning",
     "ESP-hosted SoftAP configuration. No laptop, no cable.",
     "Any technician with a phone can deploy."),
]
for i, (icon, title, body, kicker) in enumerate(solutions_row2):
    sx = Inches(0.6) + i * (sol2_w + sol2_gap)
    add_rect(s5, sx + Emu(36000), row2_y + Emu(36000), sol2_w, sol2_h, RGBColor(0xE0, 0xE3, 0xE8))
    add_rect(s5, sx, row2_y, sol2_w, sol2_h, WHITE, line=DIVIDER)
    add_rect(s5, sx, row2_y, sol2_w, Inches(0.1), AMBER)
    add_text(s5, sx + Inches(0.3), row2_y + Inches(0.25), Inches(1.0), Inches(0.6),
             icon, size=28, color=AMBER_DARK)
    add_text(s5, sx + Inches(1.4), row2_y + Inches(0.3), sol2_w - Inches(1.7), Inches(0.5),
             title, size=16, bold=True, color=NAVY)
    add_text(s5, sx + Inches(1.4), row2_y + Inches(0.8), sol2_w - Inches(1.7), Inches(0.8),
             body, size=11, color=TEXT_MID, line_spacing=1.35)
    add_rect(s5, sx, row2_y + sol2_h - Inches(0.45), sol2_w, Inches(0.45), BG_LIGHT)
    add_text(s5, sx + Inches(0.3), row2_y + sol2_h - Inches(0.4), sol2_w - Inches(0.6), Inches(0.35),
             kicker, size=11, bold=True, color=AMBER_DARK, anchor=MSO_ANCHOR.MIDDLE)

page_footer(s5, "FluxGen · v1.x Delivered", "Slide 5 of 10")

add_notes(s5, """This is where you earn credibility. Deliver it with quiet confidence — these are facts, not promises.

'Here is what v1.x delivers today, in production.

Top row — the reliability stack.
Dual connectivity gives us WiFi and cellular with auto failover.
The SD-card cache means telemetry is never lost — it replays chronologically the moment the network is back.
Auto self-healing power-cycles the modem on MQTT drops, with cooldown logic so we don't thrash.

Bottom row — the operational stack.
Remote OTA via Azure Device Twin means we push firmware and configuration from the cloud.
Web-based commissioning means any technician with a phone can deploy a gateway — no laptop.

This is not roadmap. This is shipping.'

Timing: 75 seconds. Make eye contact with the Ops head when you say 'zero data loss'.""")


# =========================================================
#  SLIDE 6 — ARCHITECTURE
# =========================================================
s6 = prs.slides.add_slide(BLANK)
page_header(s6, "How It Works", "Architecture at a glance")

arch_y = Inches(2.1)
box_h = Inches(3.5)

# FIELD box
field_x = Inches(0.6); field_w = Inches(3.2)
add_rounded(s6, field_x, arch_y, field_w, box_h, RGBColor(0xEF, 0xF7, 0xF6), line=TEAL)
add_text(s6, field_x + Inches(0.3), arch_y + Inches(0.25), field_w - Inches(0.6), Inches(0.4),
         "FIELD", size=11, bold=True, color=TEAL_DARK)
add_text(s6, field_x + Inches(0.3), arch_y + Inches(0.6), field_w - Inches(0.6), Inches(0.5),
         "Sensors", size=17, bold=True, color=NAVY)
add_bullets(s6, field_x + Inches(0.3), arch_y + Inches(1.15), field_w - Inches(0.6), Inches(2.3), [
    "Flow meters (Modbus RTU)",
    "Water quality analyzers",
    "TDS, pH, DO, temperature",
    "RS485 daisy-chain",
    "Up to 10 sensors / gateway",
], size=11, color=TEXT_MID, bullet_color=TEAL, line_spacing=1.4)

# GATEWAY box — biggest, highlighted
gw_x = Inches(4.05); gw_w = Inches(5.3)
add_rounded(s6, gw_x, arch_y, gw_w, box_h, NAVY)
add_text(s6, gw_x + Inches(0.3), arch_y + Inches(0.25), gw_w - Inches(0.6), Inches(0.4),
         "EDGE GATEWAY", size=11, bold=True, color=TEAL)
add_text(s6, gw_x + Inches(0.3), arch_y + Inches(0.6), gw_w - Inches(0.6), Inches(0.5),
         "FluxGen ESP32 Gateway", size=17, bold=True, color=WHITE)

# internal component tiles
int_tiles = [
    ("RS485 · Modbus RTU", "Configurable byte order, CRC-safe"),
    ("Dual Network Stack", "WiFi + Cellular (A7670C PPP)"),
    ("SD Offline Cache", "FAT32, chronological replay"),
    ("DS3231 RTC", "Accurate timestamps without network"),
]
tile_y = arch_y + Inches(1.2)
for i, (t, sub) in enumerate(int_tiles):
    ty = tile_y + i * Inches(0.5)
    add_rect(s6, gw_x + Inches(0.3), ty, gw_w - Inches(0.6), Inches(0.44),
             RGBColor(0x13, 0x32, 0x52))
    add_text(s6, gw_x + Inches(0.5), ty + Inches(0.05), gw_w - Inches(1.0), Inches(0.22),
             t, size=10.5, bold=True, color=WHITE)
    add_text(s6, gw_x + Inches(0.5), ty + Inches(0.22), gw_w - Inches(1.0), Inches(0.22),
             sub, size=9, color=TEAL)

# CLOUD box
cloud_x = Inches(9.55); cloud_w = Inches(3.2)
add_rounded(s6, cloud_x, arch_y, cloud_w, box_h, RGBColor(0xFD, 0xF5, 0xEC), line=AMBER)
add_text(s6, cloud_x + Inches(0.3), arch_y + Inches(0.25), cloud_w - Inches(0.6), Inches(0.4),
         "CLOUD", size=11, bold=True, color=AMBER_DARK)
add_text(s6, cloud_x + Inches(0.3), arch_y + Inches(0.6), cloud_w - Inches(0.6), Inches(0.5),
         "Azure IoT Hub", size=17, bold=True, color=NAVY)
add_bullets(s6, cloud_x + Inches(0.3), arch_y + Inches(1.15), cloud_w - Inches(0.6), Inches(2.3), [
    "MQTT · TLS secured",
    "Device Twin for config",
    "Telemetry ingestion",
    "OTA trigger channel",
    "Fleet-ready hub",
], size=11, color=TEXT_MID, bullet_color=AMBER, line_spacing=1.4)

# Arrows between boxes
def connector(slide, x1, y1, x2, y2, color):
    arrow = slide.shapes.add_connector(3, x1, y1, x2, y2)  # straight connector
    arrow.line.color.rgb = color
    arrow.line.width = Pt(2.5)
    # add arrow end
    lineEl = arrow.line._get_or_add_ln()
    tailEnd = etree.SubElement(lineEl, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    return arrow

connector(s6, field_x + field_w, arch_y + box_h/2, gw_x, arch_y + box_h/2, TEAL_DARK)
connector(s6, gw_x + gw_w, arch_y + box_h/2, cloud_x, arch_y + box_h/2, AMBER_DARK)

# Flow labels
add_text(s6, Inches(3.8), arch_y + box_h/2 + Inches(0.1), Inches(0.3), Inches(0.3),
         "RS485", size=8.5, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
add_text(s6, Inches(9.3), arch_y + box_h/2 + Inches(0.1), Inches(0.3), Inches(0.3),
         "MQTT", size=8.5, bold=True, color=AMBER_DARK, align=PP_ALIGN.CENTER)

# Bottom callout bar
add_rounded(s6, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.75), BG_LIGHT, line=DIVIDER)
add_text(s6, Inches(0.9), Inches(6.1), Inches(11.6), Inches(0.3),
         "DESIGN PHILOSOPHY", size=10, bold=True, color=TEAL_DARK)
add_text(s6, Inches(0.9), Inches(6.35), Inches(11.6), Inches(0.4),
         "Offline-first. Every layer survives the one below it failing.",
         size=14, bold=True, color=NAVY)

page_footer(s6, "FluxGen · Architecture", "Slide 6 of 10")

add_notes(s6, """Draw with your hands — field on the left, cloud on the right, gateway in the middle.

'The architecture has three layers.

On the left — the field. Up to ten sensors, flow meters and water quality, talking to us over RS485 Modbus.

In the middle — our gateway. This is where the intelligence lives. RS485 in, dual network stack, SD card cache, real-time clock for timestamp accuracy. Everything inside this box is designed around one principle — offline-first. Every layer survives the one below it failing.

On the right — Azure IoT Hub. MQTT over TLS, Device Twin for remote configuration, telemetry in, OTA out.

The callout at the bottom is the key idea — offline-first. That is why v1.x doesn't lose data.'

Timing: 60 seconds.""")


# =========================================================
#  SLIDE 7 — BUSINESS IMPACT / ROI
# =========================================================
s7 = prs.slides.add_slide(BLANK)
page_header(s7, "Business Impact", "What v1.x delivers, per site, per year")

# Left — chart
chart_data = CategoryChartData()
chart_data.categories = ["Imported PLC\n(typical)", "FluxGen v1.x"]
chart_data.add_series("Annual cost per site (₹ ’000)", (85, 32))

chart = s7.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.6), Inches(1.8), Inches(5.5), Inches(4.8),
    chart_data,
).chart
chart.has_title = True
chart.chart_title.text_frame.text = "Annual cost per site — ₹ thousands"
for run in chart.chart_title.text_frame.paragraphs[0].runs:
    run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = NAVY; run.font.name = FONT
chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True
dl = plot.data_labels
dl.font.size = Pt(11); dl.font.bold = True; dl.font.color.rgb = NAVY
dl.number_format = "0"
dl.position = XL_LABEL_POSITION.OUTSIDE_END
# color each bar
series = plot.series[0]
# Apply varied colors via per-point fill
from pptx.oxml.ns import nsmap
ser_xml = series._element
# Brute force: set two data points with different colors
pts_xml = etree.SubElement(ser_xml, qn("c:dPt"))
idx = etree.SubElement(pts_xml, qn("c:idx")); idx.set("val", "0")
invertIfNeg = etree.SubElement(pts_xml, qn("c:invertIfNegative")); invertIfNeg.set("val", "0")
bubble3d = etree.SubElement(pts_xml, qn("c:bubble3D")); bubble3d.set("val", "0")
spPr = etree.SubElement(pts_xml, qn("c:spPr"))
fill = etree.SubElement(spPr, qn("a:solidFill"))
srgb = etree.SubElement(fill, qn("a:srgbClr")); srgb.set("val", "F4A261")
pts_xml2 = etree.SubElement(ser_xml, qn("c:dPt"))
idx2 = etree.SubElement(pts_xml2, qn("c:idx")); idx2.set("val", "1")
invertIfNeg2 = etree.SubElement(pts_xml2, qn("c:invertIfNegative")); invertIfNeg2.set("val", "0")
bubble3d2 = etree.SubElement(pts_xml2, qn("c:bubble3D")); bubble3d2.set("val", "0")
spPr2 = etree.SubElement(pts_xml2, qn("c:spPr"))
fill2 = etree.SubElement(spPr2, qn("a:solidFill"))
srgb2 = etree.SubElement(fill2, qn("a:srgbClr")); srgb2.set("val", "0A2540")

# Right — impact table (custom built)
right_x = Inches(6.4); right_w = Inches(6.4)
add_text(s7, right_x, Inches(1.8), right_w, Inches(0.4),
         "PER-SITE IMPACT (INDICATIVE)", size=11, bold=True, color=AMBER_DARK)
add_text(s7, right_x, Inches(2.15), right_w, Inches(0.5),
         "Where the value shows up.", size=20, bold=True, color=NAVY)

rows = [
    ("Field service visits avoided", "4–6 / yr", "₹ 20,000"),
    ("Firmware update cost", "near-zero", "₹ 12,000"),
    ("Manual data collection", "eliminated", "₹ 18,000"),
    ("Compliance-risk exposure", "mitigated", "qualitative"),
    ("Device cost vs imported PLC", "50–70% lower", "CapEx saving"),
]
row_y = Inches(2.85); row_h = Inches(0.62)
# header
add_rect(s7, right_x, row_y, right_w, Inches(0.42), NAVY)
add_text(s7, right_x + Inches(0.2), row_y + Inches(0.08), Inches(3.0), Inches(0.3),
         "METRIC", size=10, bold=True, color=WHITE)
add_text(s7, right_x + Inches(3.4), row_y + Inches(0.08), Inches(1.6), Inches(0.3),
         "CHANGE", size=10, bold=True, color=WHITE)
add_text(s7, right_x + Inches(5.0), row_y + Inches(0.08), Inches(1.3), Inches(0.3),
         "VALUE / YR", size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

cur_y = row_y + Inches(0.42)
for i, (m, c, v) in enumerate(rows):
    bg = WHITE if i % 2 == 0 else BG_LIGHT
    add_rect(s7, right_x, cur_y, right_w, row_h, bg)
    add_text(s7, right_x + Inches(0.2), cur_y + Inches(0.18), Inches(3.1), Inches(0.3),
             m, size=12, color=TEXT_DARK)
    add_text(s7, right_x + Inches(3.4), cur_y + Inches(0.18), Inches(1.6), Inches(0.3),
             c, size=12, color=TEXT_MID)
    add_text(s7, right_x + Inches(5.0), cur_y + Inches(0.18), Inches(1.3), Inches(0.3),
             v, size=12, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    cur_y += row_h

# total callout
add_rounded(s7, right_x, cur_y + Inches(0.1), right_w, Inches(0.65), AMBER)
add_text(s7, right_x + Inches(0.3), cur_y + Inches(0.18), right_w - Inches(0.6), Inches(0.5),
         "Indicative direct savings: ~₹ 50,000 per site, per year.",
         size=14, bold=True, color=NAVY_DARK, anchor=MSO_ANCHOR.MIDDLE)

# Footnote
add_text(s7, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
         "Note: Indicative figures for illustration. Actuals vary by site profile and customer SLA.",
         size=9, color=TEXT_LIGHT)

page_footer(s7, "FluxGen · Business Impact", "Slide 7 of 10")

add_notes(s7, """This is the CFO slide. Lead with the number, not the methodology.

'Let me put a number on this. Per site, per year, v1.x delivers roughly fifty thousand rupees in direct savings — before we count the strategic value of compliance and customer trust.

Four to six field service visits avoided. Firmware updates at near-zero cost. Manual data collection eliminated. And our BOM sits at fifty to seventy percent below an imported PLC doing the same job.

The chart on the left is the takeaway — the imported alternative costs about eighty-five thousand rupees per site per year; our gateway does it for about thirty-two.

These are indicative numbers — I have the build-up if you want to go deeper.'

If CFO pushes for specifics: 'These are based on typical Indian STP/ETP deployments — I'll pull exact numbers from our top three customer sites after this meeting.' Don't fake precision you don't have.

Timing: 75 seconds.""")


# =========================================================
#  SLIDE 8 — v2.0 ROADMAP
# =========================================================
s8 = prs.slides.add_slide(BLANK)
page_header(s8, "v2.0 Proposal", "From device reliability to fleet intelligence")

# Intro
add_text(s8, Inches(0.6), Inches(1.55), Inches(12), Inches(0.4),
         "Five features. Each closes a concrete gap between where v1.x sits and a true platform.",
         size=13, color=TEXT_MID)

# 5 rows, compact
features = [
    ("1", "Fleet Management Console",
     "Single pane of glass — all gateways, bulk OTA, health alerts.",
     "Operational leverage", NAVY),
    ("2", "Enterprise Security Hardening",
     "Signed OTA, TLS cert pinning, secure boot, audit trails.",
     "Unlocks tenders", AMBER),
    ("3", "Edge Analytics",
     "On-device anomaly detection — alert before the cloud sees it.",
     "Differentiator", TEAL),
    ("4", "Field Technician Mobile App",
     "Commission, diagnose, and recover gateways from a phone.",
     "Ops efficiency", TEAL_DARK),
    ("5", "Modular Firmware Refactor",
     "Split 609KB monolith → faster features, fewer regressions.",
     "4× delivery speed", AMBER_DARK),
]
row_start_y = Inches(2.1); row_height = Inches(0.9); row_gap = Inches(0.06)
for i, (num, name, desc, tag, col) in enumerate(features):
    ry = row_start_y + i * (row_height + row_gap)
    add_rect(s8, Inches(0.6), ry, Inches(12.15), row_height, WHITE, line=DIVIDER)
    # left color bar
    add_rect(s8, Inches(0.6), ry, Inches(0.15), row_height, col)
    # number
    add_text(s8, Inches(0.95), ry + Inches(0.2), Inches(0.6), Inches(0.5),
             num, size=28, bold=True, color=col)
    # name + desc
    add_text(s8, Inches(1.75), ry + Inches(0.12), Inches(7.5), Inches(0.4),
             name, size=15, bold=True, color=NAVY)
    add_text(s8, Inches(1.75), ry + Inches(0.48), Inches(7.5), Inches(0.4),
             desc, size=11.5, color=TEXT_MID)
    # tag pill
    pill(s8, Inches(9.6), ry + Inches(0.3), Inches(3.0), Inches(0.35),
         tag.upper(), fill=col, color=WHITE, size=10)

# Timeline footer
tl_y = Inches(6.85)
add_text(s8, Inches(0.6), tl_y, Inches(12), Inches(0.3),
         "TIMELINE: 16 weeks · 4 milestones · go/no-go gate after each",
         size=10, bold=True, color=TEXT_MID)

page_footer(s8, "FluxGen · v2.0 Roadmap", "Slide 8 of 10")

add_notes(s8, """This is the vision slide. Read the headlines, don't read the slide word-for-word.

'v2.0 is five features. Each one closes a specific gap.

ONE — Fleet Management Console. We have a great device today. We don't have a way to see all of them at once. This is the single biggest ops leverage play in the list.

TWO — Enterprise Security Hardening. Signed OTA, TLS cert pinning, secure boot. This is not optional — it is the gate to government and enterprise tenders.

THREE — Edge Analytics. On-device anomaly detection. This is where we stop looking like a gateway vendor and start looking like a platform.

FOUR — Field Technician Mobile App. Commission from a phone. Halves commissioning time.

FIVE — Modular Firmware Refactor. The current web-config file is six hundred kilobytes in one file. Splitting it unlocks about four times the feature-delivery speed.

All of this in sixteen weeks, with a go/no-go gate after every milestone.'

Timing: 80 seconds — don't rush the security tenders line, it matters most to the CEO.""")


# =========================================================
#  SLIDE 9 — THE ASK + RISKS
# =========================================================
s9 = prs.slides.add_slide(BLANK)
page_header(s9, "The Ask", "What I need to deliver v2.0")

# Left column — the ask
left_x = Inches(0.6); left_w = Inches(6.1)
add_rounded(s9, left_x, Inches(1.7), left_w, Inches(5.05), NAVY)
add_text(s9, left_x + Inches(0.4), Inches(1.9), left_w - Inches(0.8), Inches(0.4),
         "INVESTMENT", size=11, bold=True, color=TEAL)
add_text(s9, left_x + Inches(0.4), Inches(2.25), left_w - Inches(0.8), Inches(0.6),
         "What I'm asking for.", size=22, bold=True, color=WHITE)

ask_items = [
    ("Team", "1 Firmware engineer + 1 Cloud / Full-stack engineer"),
    ("Timeline", "16 weeks · phased delivery · go/no-go gates"),
    ("Budget", "Salary + ~₹1.5L cloud/tooling (indicative)"),
    ("My role", "IoT Platform Tech Lead — device-to-cloud ownership"),
    ("Deliverables", "Fleet console · Secure OTA · Mobile app · Edge analytics · Refactored firmware"),
]
cur_y = Inches(3.0)
for label, value in ask_items:
    add_text(s9, left_x + Inches(0.4), cur_y, Inches(1.5), Inches(0.3),
             label.upper(), size=10, bold=True, color=TEAL)
    add_text(s9, left_x + Inches(0.4), cur_y + Inches(0.28), left_w - Inches(0.8), Inches(0.5),
             value, size=13, color=WHITE, line_spacing=1.3)
    cur_y += Inches(0.72)

# Right column — risks table
right_x = Inches(6.9); right_w = Inches(5.85)
add_text(s9, right_x, Inches(1.85), right_w, Inches(0.4),
         "RISKS · MITIGATION", size=11, bold=True, color=AMBER_DARK)
add_text(s9, right_x, Inches(2.2), right_w, Inches(0.5),
         "What could go wrong, what we do about it.", size=16, bold=True, color=NAVY)

risks = [
    ("Key-person dependency", "Refactor + docs in Month 1 before new features."),
    ("Scope creep", "Fixed 16-week window · gated milestones."),
    ("Security gaps in current build", "Security hardening is Milestone 2 — before edge analytics."),
    ("Customer impact during rollout", "Staged OTA · canary devices first · rollback built in."),
    ("Cloud cost drift", "Monthly Azure review · feature flags for analytics."),
]
tr_y = Inches(2.95)
add_rect(s9, right_x, tr_y, right_w, Inches(0.4), NAVY)
add_text(s9, right_x + Inches(0.2), tr_y + Inches(0.08), Inches(2.3), Inches(0.3),
         "RISK", size=10, bold=True, color=WHITE)
add_text(s9, right_x + Inches(2.55), tr_y + Inches(0.08), Inches(3.0), Inches(0.3),
         "MITIGATION", size=10, bold=True, color=WHITE)

ry_cur = tr_y + Inches(0.4)
for i, (risk, mit) in enumerate(risks):
    bg = WHITE if i % 2 == 0 else BG_LIGHT
    add_rect(s9, right_x, ry_cur, right_w, Inches(0.63), bg)
    add_text(s9, right_x + Inches(0.2), ry_cur + Inches(0.12), Inches(2.3), Inches(0.5),
             risk, size=11, bold=True, color=NAVY, line_spacing=1.25)
    add_text(s9, right_x + Inches(2.55), ry_cur + Inches(0.12), Inches(3.1), Inches(0.5),
             mit, size=11, color=TEXT_MID, line_spacing=1.25)
    ry_cur += Inches(0.63)

page_footer(s9, "FluxGen · The Ask", "Slide 9 of 10")

add_notes(s9, """This is THE slide. Don't hedge.

'Here is what I'm asking for.

Team — two engineers. One firmware, one cloud. With them I can deliver all five v2.0 features.

Timeline — sixteen weeks, phased, go/no-go gates. If we hit Milestone 1 and the business wants to reprioritise, nothing is wasted.

Budget — their salaries plus about a lakh and a half for cloud and tooling. Indicative.

My role — I want the mandate as IoT Platform Tech Lead. Ownership of the stack end to end — device to cloud. I've been effectively doing this; I'm asking for it to be formal so I can scale the team.

Deliverables are listed.

On the right — the risks I've thought about and how I handle each. Key-person risk, scope creep, security, customer impact, cost drift. Nothing here surprises me.'

Expect pushback. Prep answers: see companion doc. Timing: 90 seconds — you can take more time here.""")


# =========================================================
#  SLIDE 10 — CLOSING
# =========================================================
s10 = prs.slides.add_slide(BLANK)
# Full-bleed navy
add_rect(s10, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
# Teal corner accent
add_rect(s10, Inches(0), Inches(0), Inches(0.25), SLIDE_H, TEAL)
add_rect(s10, SLIDE_W - Inches(0.25), Inches(0), Inches(0.25), SLIDE_H, AMBER)

add_text(s10, Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.4),
         "CLOSING", size=12, bold=True, color=TEAL)

add_text(s10, Inches(1.5), Inches(1.7), Inches(10.3), Inches(2.4),
         "v1.x proved we can build\nfield-grade reliability\nthat global brands charge 3× for.",
         size=36, bold=True, color=WHITE, line_spacing=1.15)

add_text(s10, Inches(1.5), Inches(4.35), Inches(10.3), Inches(1.2),
         "v2.0 turns that foundation into FluxGen's platform advantage.",
         size=22, color=TEAL, line_spacing=1.2)

# CTA pill
add_rounded(s10, Inches(1.5), Inches(5.7), Inches(4.0), Inches(0.7), AMBER)
add_text(s10, Inches(1.5), Inches(5.7), Inches(4.0), Inches(0.7),
         "Greenlight v2.0 →", size=18, bold=True, color=NAVY_DARK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s10, Inches(1.5), Inches(6.65), Inches(10.3), Inches(0.3),
         "Y Chaithanya Reddy   ·   FluxGen   ·   IoT Platform",
         size=11, color=TEXT_LIGHT)

add_notes(s10, """Don't read the slide. Deliver this from memory, quietly, with weight.

'To close —

v1.x proved we can build field-grade reliability that global brands charge three times as much for.

v2.0 turns that foundation into FluxGen's platform advantage.

I'm asking for the mandate and the team to make that happen.

Thank you. Happy to take questions.'

Then stop talking. Pause. Let the silence work. Don't fill it.""")


# ---------- SAVE ----------
out_path = r"C:\Users\chath\Documents\Python code\modbus_iot_gateway\presentation\FluxGen_IoT_Gateway_v2.0_Pitch.pptx"
prs.save(out_path)
print(f"Generated: {out_path}")
print(f"Slides: {len(prs.slides)}")
