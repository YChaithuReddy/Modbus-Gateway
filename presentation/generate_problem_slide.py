"""
Standalone generator — produces a single .pptx file containing ONLY the
Problem Statement slide in a simple paragraph-style layout.
Pure typography. No tiles. No icons. No grid.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ---------- PALETTE ----------
NAVY       = RGBColor(0x0A, 0x25, 0x40)
TEAL       = RGBColor(0x00, 0xC2, 0xB8)
TEAL_DARK  = RGBColor(0x00, 0x8C, 0x85)
ALERT      = RGBColor(0xE6, 0x39, 0x46)
BG_LIGHT   = RGBColor(0xF8, 0xF9, 0xFB)
DIVIDER    = RGBColor(0xE5, 0xE8, 0xEC)
TEXT_DARK  = RGBColor(0x2A, 0x2E, 0x37)
TEXT_MID   = RGBColor(0x5B, 0x64, 0x72)
TEXT_LIGHT = RGBColor(0x95, 0x9C, 0xA8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------- HELPERS ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color=DIVIDER, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


# ---------- SLIDE ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)

# --- Header (kicker + title + divider) ---
add_text(s, Inches(0.6), Inches(0.35), Inches(11), Inches(0.3),
         "THE PROBLEM", size=11, bold=True, color=ALERT)
add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
         "What was broken before Aquagen", size=30, bold=True, color=NAVY)
add_line(s, Inches(0.6), Inches(1.35), Inches(12.7), Inches(1.35), DIVIDER, 1.0)

# --- Left vertical accent bar (subtle visual anchor) ---
add_rect(s, Inches(0.6), Inches(1.8), Inches(0.1), Inches(4.5), TEAL)

# --- Big statement (2 lines) ---
add_text(s, Inches(1.0), Inches(1.75), Inches(11.6), Inches(1.4),
         "The device was losing data in the field —\nand locking out non-developers.",
         size=28, bold=True, color=NAVY, line_spacing=1.25)

# --- Paragraphs (narrative) ---
paragraph_x = Inches(1.0)
paragraph_w = Inches(11.6)

# Paragraph 1 — WiFi problem
add_text(s, paragraph_x, Inches(3.25), paragraph_w, Inches(1.0),
         "Our ESP32 IoT device could only connect over WiFi. At many installation sites the WiFi was "
         "unstable or unavailable altogether. Whenever the network dropped, telemetry from that window "
         "was lost permanently — no local buffer, no replay, no recovery.",
         size=14, color=TEXT_DARK, line_spacing=1.55)

# Paragraph 2 — Config problem
add_text(s, paragraph_x, Inches(4.45), paragraph_w, Inches(1.0),
         "At the same time, every configuration change — however small — required opening the Arduino "
         "source code, editing it manually, recompiling, and reflashing the firmware into the ESP32. "
         "Every update became a developer-only task, locking out non-technical users entirely.",
         size=14, color=TEXT_DARK, line_spacing=1.55)

# Paragraph 3 — Business impact
add_text(s, paragraph_x, Inches(5.65), paragraph_w, Inches(1.0),
         "The result was data loss, compliance gaps in customer reports, constant truck rolls for field "
         "visits, and a device that could not be deployed at scale without engineering involvement.",
         size=14, color=TEXT_DARK, line_spacing=1.55)

# --- Bottom pull quote (muted italic) ---
add_text(s, Inches(1.0), Inches(6.9), Inches(11.6), Inches(0.4),
         "Two failures, one device: connectivity and configuration.",
         size=13, italic=True, color=TEAL_DARK, line_spacing=1.2)

# --- Speaker notes ---
notes = s.notes_slide.notes_text_frame
notes.text = """Deliver this as a narrative — read the paragraphs smoothly, don't rush.

'Before we built Aquagen, the device had a simple but damaging pair of problems.

It was losing data in the field, and it was locking out anyone who wasn't a developer.

The first problem was connectivity. Our ESP32 device could only connect over WiFi. At many installation sites, the WiFi was unstable or simply not available. Whenever the network dropped, telemetry from that window was lost permanently — the device had no local buffer, no replay, no recovery.

The second problem was configuration. Every change — however small — required opening the Arduino source code, editing it manually, recompiling, and reflashing the firmware into the ESP32. Every update became a developer-only task. Non-technical users were locked out entirely.

The result — data loss, compliance gaps in customer reports, constant truck rolls for field visits, and a device that could not be deployed at scale without engineering involvement.

Two failures, one device: connectivity and configuration. This is the mountain we set out to move.'

Timing: 75 seconds."""


# ---------- SAVE ----------
OUT = Path(__file__).parent / "Problem_Statement.pptx"
prs.save(OUT)
print(f"Generated: {OUT}")
print(f"Size: {OUT.stat().st_size // 1024} KB")
