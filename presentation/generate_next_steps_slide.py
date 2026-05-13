"""
Standalone generator — produces a single .pptx file containing ONLY the
Next Steps slide. Clean typography, paragraph-style.
Copy-paste friendly into any deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ---------- PALETTE ----------
NAVY       = RGBColor(0x0A, 0x25, 0x40)
TEAL       = RGBColor(0x00, 0xC2, 0xB8)
TEAL_DARK  = RGBColor(0x00, 0x8C, 0x85)
AMBER      = RGBColor(0xF4, 0xA2, 0x61)
AMBER_DARK = RGBColor(0xD4, 0x84, 0x41)
BG_LIGHT   = RGBColor(0xF8, 0xF9, 0xFB)
DIVIDER    = RGBColor(0xE5, 0xE8, 0xEC)
TEXT_DARK  = RGBColor(0x2A, 0x2E, 0x37)
TEXT_MID   = RGBColor(0x5B, 0x64, 0x72)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------- HELPERS ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_rounded(slide, x, y, w, h, fill, line=None, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color=DIVIDER, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


# ---------- SLIDE ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)

# --- Header ---
add_text(s, Inches(0.6), Inches(0.35), Inches(11), Inches(0.3),
         "NEXT STEPS", size=11, bold=True, color=TEAL_DARK)
add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
         "Where we go from here", size=30, bold=True, color=NAVY)
add_line(s, Inches(0.6), Inches(1.35), Inches(12.7), Inches(1.35), DIVIDER, 1.0)

# --- Intro line ---
add_text(s, Inches(0.6), Inches(1.65), Inches(12.15), Inches(0.4),
         "Two steps to take the device from working prototype to production-ready product.",
         size=14, color=TEXT_MID, line_spacing=1.3)

# ============================================================
# STEP 01 — Design the integrated PCB
# ============================================================
step1_y = Inches(2.4)

# Big "01" teal number
add_text(s, Inches(0.6), step1_y, Inches(1.6), Inches(1.3),
         "01", size=78, bold=True, color=TEAL, line_spacing=1.0)

# Step title
add_text(s, Inches(2.2), step1_y + Inches(0.05), Inches(10.5), Inches(0.55),
         "Design the integrated PCB",
         size=24, bold=True, color=NAVY, line_spacing=1.1)

# Step description
add_text(s, Inches(2.2), step1_y + Inches(0.65), Inches(10.5), Inches(0.9),
         "Consolidate every new feature — SIM module, SD card backup, RTC, and web-server support — "
         "onto a single custom-designed PCB. This becomes our production board and replaces the current "
         "multi-module prototype.",
         size=13, color=TEXT_DARK, line_spacing=1.55)

# Owner tag — "Led by Mallikarjun"
owner_y = step1_y + Inches(1.55)
add_rounded(s, Inches(2.2), owner_y, Inches(3.4), Inches(0.35), TEAL, radius=0.5)
add_text(s, Inches(2.35), owner_y + Inches(0.03), Inches(3.2), Inches(0.3),
         "Led by Mallikarjun  ·  In progress",
         size=10.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# STEP 02 — Scale up production
# ============================================================
step2_y = Inches(4.65)

# Big "02" amber number
add_text(s, Inches(0.6), step2_y, Inches(1.6), Inches(1.3),
         "02", size=78, bold=True, color=AMBER, line_spacing=1.0)

# Step title
add_text(s, Inches(2.2), step2_y + Inches(0.05), Inches(10.5), Inches(0.55),
         "Scale up production",
         size=24, bold=True, color=NAVY, line_spacing=1.1)

# Step description
add_text(s, Inches(2.2), step2_y + Inches(0.65), Inches(10.5), Inches(0.9),
         "Ramp up manufacturing of the integrated PCB so upcoming deployments ship with the full feature "
         "set from day one. Every new project gets a device that handles outages gracefully — zero data "
         "loss, out of the box.",
         size=13, color=TEXT_DARK, line_spacing=1.55)

# --- Bottom takeaway pill ---
takeaway_y = Inches(7.15)
add_rounded(s, Inches(0.6), takeaway_y - Inches(0.03), Inches(12.15), Inches(0.4),
            NAVY, radius=0.4)
add_text(s, Inches(0.9), takeaway_y, Inches(11.5), Inches(0.35),
         "From prototype to scale — one integrated PCB at a time.",
         size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

# --- Speaker notes ---
notes = s.notes_slide.notes_text_frame
notes.text = """Deliver this as forward-looking and confident — not as a request, as a direction.

'So what's next? Two clear steps.

STEP ONE — design the integrated PCB. Right now, the device uses a multi-module prototype. For production, we're consolidating every new feature — the SIM module, SD card, RTC, and web server hardware — onto a single custom-designed PCB. This will be our production board. Mallikarjun is leading this work, currently in progress.

STEP TWO — once the integrated PCB is ready, we scale up production. The goal is that every upcoming deployment ships with the full feature set from day one. New projects get a device that handles outages gracefully — with zero data loss, out of the box.

From working prototype to scalable product. One integrated PCB at a time.'

Timing: 60 seconds."""


# ---------- SAVE ----------
OUT = Path(__file__).parent / "Next_Steps.pptx"
prs.save(OUT)
print(f"Generated: {OUT}")
print(f"Size: {OUT.stat().st_size // 1024} KB")
