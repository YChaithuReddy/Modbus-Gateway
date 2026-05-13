"""
ESP32 IoT Device Enhancement — Project Showcase Deck
9 slides · 10-minute office presentation
Structure follows the user's 9-section brief verbatim.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

# ---------- PALETTE ----------
NAVY       = RGBColor(0x0A, 0x25, 0x40)
NAVY_DARK  = RGBColor(0x06, 0x17, 0x2A)
TEAL       = RGBColor(0x00, 0xC2, 0xB8)
TEAL_DARK  = RGBColor(0x00, 0x8C, 0x85)
AMBER      = RGBColor(0xF4, 0xA2, 0x61)
AMBER_DARK = RGBColor(0xD4, 0x84, 0x41)
BG_LIGHT   = RGBColor(0xF8, 0xF9, 0xFB)
DIVIDER    = RGBColor(0xE5, 0xE8, 0xEC)
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MID   = RGBColor(0x5B, 0x64, 0x72)
TEXT_LIGHT = RGBColor(0x95, 0x9C, 0xA8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ALERT      = RGBColor(0xE6, 0x39, 0x46)
SUCCESS    = RGBColor(0x2A, 0x9D, 0x8F)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

HERE = Path(__file__).parent
ASSETS = HERE / "assets"
ICONS = ASSETS / "icons"


# ============================================================
# HELPERS (unchanged)
# ============================================================
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_rounded(slide, x, y, w, h, fill, line=None, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_DARK,
             font=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT_DARK,
                bullet_color=None, font=FONT, line_spacing=1.35):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        r0 = p.add_run(); r0.text = "●  "
        r0.font.name = font; r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color or color; r0.font.bold = True
        if isinstance(item, tuple):
            bold_part, rest = item
            r1 = p.add_run(); r1.text = bold_part
            r1.font.name = font; r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = color
            r2 = p.add_run(); r2.text = rest
            r2.font.name = font; r2.font.size = Pt(size); r2.font.color.rgb = color
        else:
            r1 = p.add_run(); r1.text = item
            r1.font.name = font; r1.font.size = Pt(size); r1.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color=DIVIDER, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    elif w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    elif h:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    return slide.shapes.add_picture(str(path), x, y)


def page_header(slide, kicker, title, *, kicker_color=TEAL):
    add_text(slide, Inches(0.6), Inches(0.35), Inches(11), Inches(0.3),
             kicker.upper(), size=11, bold=True, color=kicker_color)
    add_text(slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    add_line(slide, Inches(0.6), Inches(1.35), Inches(12.7), Inches(1.35), DIVIDER, 1.0)


def page_footer(slide, text_left="ESP32 IoT Enhancement", text_right=""):
    add_text(slide, Inches(0.6), Inches(7.1), Inches(6), Inches(0.3),
             text_left, size=9, color=TEXT_LIGHT)
    if text_right:
        add_text(slide, Inches(7.3), Inches(7.1), Inches(5.4), Inches(0.3),
                 text_right, size=9, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def pill(slide, x, y, w, h, text, *, fill=TEAL, color=WHITE, size=10):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    tf = s.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color


# ============================================================
# TRANSITION INJECTION
# ============================================================
NS_P15 = "http://schemas.microsoft.com/office/powerpoint/2012/main"


def add_slide_transition(slide, transition_type="fade"):
    sld = slide.element
    for t in sld.findall(qn("p:transition")):
        sld.remove(t)
    trans = etree.SubElement(sld, qn("p:transition"))
    trans.set("spd", "med")
    trans.set("advClick", "1")
    if transition_type == "fade":
        etree.SubElement(trans, qn("p:fade"))
    return trans


# ============================================================
#  SLIDE 1 — TITLE
# ============================================================
s1 = prs.slides.add_slide(BLANK)
add_image(s1, ASSETS / "bg_title.png", Inches(0), Inches(0), SLIDE_W, SLIDE_H)

add_text(s1, Inches(0.9), Inches(0.85), Inches(8), Inches(0.4),
         "PROJECT SHOWCASE", size=12, bold=True, color=TEAL)

add_text(s1, Inches(0.9), Inches(1.9), Inches(12), Inches(1.2),
         "ESP32 IoT Device", size=62, bold=True, color=WHITE, line_spacing=1.0)
add_text(s1, Inches(0.9), Inches(2.95), Inches(12), Inches(1.2),
         "Enhancement", size=62, bold=True, color=WHITE, line_spacing=1.0)

add_rect(s1, Inches(0.95), Inches(4.25), Inches(1.2), Inches(0.05), TEAL)

add_text(s1, Inches(0.9), Inches(4.45), Inches(11.8), Inches(0.5),
         "SIM Connectivity   ·   SD Card Backup   ·   Web-Based Configuration",
         size=20, color=TEAL)

add_text(s1, Inches(0.9), Inches(5.05), Inches(11.8), Inches(0.4),
         "From WiFi-dependent to smart hybrid IoT solution",
         size=14, color=TEXT_LIGHT)

# Presenter panel
panel = add_rounded(s1, Inches(0.9), Inches(5.85), Inches(5.5), Inches(1.25),
                    RGBColor(0x0E, 0x2C, 0x4C), radius=0.12)
panel.line.color.rgb = RGBColor(0x1F, 0x4A, 0x78)
panel.line.width = Pt(0.75)

add_text(s1, Inches(1.15), Inches(6.0), Inches(5.0), Inches(0.3),
         "PRESENTED BY", size=10, bold=True, color=TEAL)
add_text(s1, Inches(1.15), Inches(6.3), Inches(5.0), Inches(0.35),
         "Y Chaithanya Reddy", size=17, bold=True, color=WHITE)
add_text(s1, Inches(1.15), Inches(6.65), Inches(5.0), Inches(0.3),
         "IoT Platform  ·  Firmware & Cloud  ·  April 2026", size=10, color=TEXT_LIGHT)

add_notes(s1, """Open strong — eye contact before first word.

'Good morning. In the next ten minutes, I'll walk you through a project that transformed our ESP32 IoT device — from a WiFi-dependent system into a smart hybrid solution that handles network outages, stores data safely offline, and lets anyone configure it through a web browser.

Three upgrades, one simple story. Let's go through it.'

Pause 2 seconds. Click.""")


# ============================================================
#  SLIDE 2 — PREVIOUS SYSTEM CHALLENGES
# ============================================================
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)
page_header(s2, "Previous System Challenges", "What was breaking before the upgrade")

# LEFT PANEL — field reality
left_x = Inches(0.6); left_w = Inches(6.0)
add_rounded(s2, left_x, Inches(1.75), left_w, Inches(5.1),
            RGBColor(0xFD, 0xF3, 0xF3), radius=0.04)
add_text(s2, left_x + Inches(0.4), Inches(1.95), left_w - Inches(0.8), Inches(0.35),
         "FIELD REALITY", size=11, bold=True, color=ALERT)
add_text(s2, left_x + Inches(0.4), Inches(2.3), left_w - Inches(0.8), Inches(0.6),
         "WiFi-only. Code-only.", size=22, bold=True, color=NAVY, line_spacing=1.1)
add_bullets(s2, left_x + Inches(0.4), Inches(3.1), left_w - Inches(0.8), Inches(3.7), [
    ("Only WiFi connectivity → ", "no cellular fallback available."),
    ("Unstable WiFi at many sites → ", "frequent network downtime."),
    ("Network drops → ", "significant data loss, no recovery."),
    ("Any config change → ", "open Arduino code, edit, recompile."),
    ("Every update → ", "reflash firmware into the ESP32 manually."),
], size=13, color=TEXT_DARK, bullet_color=ALERT, line_spacing=1.6)

# RIGHT PANEL — 2x2 impact tiles
right_x = Inches(6.9); right_w = Inches(5.85)
add_text(s2, right_x, Inches(1.85), right_w, Inches(0.4),
         "IMPACT", size=11, bold=True, color=AMBER_DARK)
add_text(s2, right_x, Inches(2.2), right_w, Inches(0.6),
         "The cost of the old way.", size=22, bold=True, color=NAVY, line_spacing=1.1)

tile_w = Inches(2.85); tile_h = Inches(1.85); tile_gap = Inches(0.15)
impacts = [
    ("Data loss",
     "Every network outage lost telemetry permanently.",
     ALERT, "p_dataloss.png"),
    ("Technical dependency",
     "Every change required a developer and recompile.",
     AMBER_DARK, "p_laptop.png"),
    ("Slow rollout",
     "Each site needed code edits + manual reflash.",
     ALERT, "p_truck.png"),
    ("Poor user experience",
     "Non-technical users locked out of configuration.",
     AMBER_DARK, "p_eye.png"),
]
for i, (title, body, col, icon_file) in enumerate(impacts):
    col_idx = i % 2; row_idx = i // 2
    tx = right_x + col_idx * (tile_w + tile_gap)
    ty = Inches(3.0) + row_idx * (tile_h + tile_gap)
    add_rect(s2, tx, ty, tile_w, tile_h, WHITE, line=DIVIDER)
    add_rect(s2, tx, ty, Inches(0.08), tile_h, col)
    add_image(s2, ICONS / icon_file, tx + Inches(0.25), ty + Inches(0.2),
              w=Inches(0.55), h=Inches(0.55))
    add_text(s2, tx + Inches(0.9), ty + Inches(0.25), tile_w - Inches(1.0), Inches(0.4),
             title, size=13, bold=True, color=NAVY)
    add_text(s2, tx + Inches(0.25), ty + Inches(0.9), tile_w - Inches(0.4), tile_h - Inches(1.0),
             body, size=11, color=TEXT_MID, line_spacing=1.35)

page_footer(s2, "ESP32 IoT Enhancement · Previous Challenges", "Slide 2 of 11")

add_notes(s2, """'Before this upgrade, the device had two big problems.

ONE — connectivity. It only supported WiFi. At many installation sites, WiFi was unstable or simply unavailable. Whenever the network dropped, data was lost permanently. There was no local buffer, no replay, no recovery.

TWO — configuration. To change any setting — even something small — we had to open the Arduino source code, modify it manually, and reflash the firmware into the ESP32. This was technical, time-consuming, and locked out non-technical users completely.

On the right — the real-world impact. Data loss. Technical dependency. Slow rollout. Poor experience for anyone who was not a developer.

This is the baseline we set out to fix.'

Timing: 75 seconds.""")


# ============================================================
#  SLIDE 3 — NEW SOLUTION OVERVIEW
# ============================================================
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)
page_header(s3, "New Solution Implemented", "Four additions that transformed the device")

# 2x2 grid of solution cards
cards_x0 = Inches(0.6)
cards_y0 = Inches(1.8)
card_w = Inches(6.05)
card_h = Inches(2.45)
card_gap_x = Inches(0.15)
card_gap_y = Inches(0.15)

solutions = [
    ("wifi.png", "SIM Module", "01", TEAL,
     "Cellular connectivity — device now uses mobile SIM network directly, reducing WiFi dependency and improving uptime in remote or unstable areas."),
    ("sdcard.png", "SD Card Backup", "02", TEAL,
     "Automatic offline storage — when network drops, SD card activates and stores every data ping locally. Zero data loss during outages."),
    ("refresh.png", "RTC Module", "03", AMBER,
     "Real-Time Clock — provides accurate date and time continuously, so every stored ping gets an exact timestamp. Reliable historical records."),
    ("globe.png", "Web Server", "04", AMBER,
     "Browser-based configuration portal — configure device settings from any browser. No Arduino editing, no firmware reflash, no developer needed."),
]

for i, (icon_f, title, num, accent, body) in enumerate(solutions):
    col_idx = i % 2; row_idx = i // 2
    cx = cards_x0 + col_idx * (card_w + card_gap_x)
    cy = cards_y0 + row_idx * (card_h + card_gap_y)
    # shadow
    add_rect(s3, cx + Emu(45000), cy + Emu(45000), card_w, card_h, RGBColor(0xDD, 0xE1, 0xE7))
    # body
    add_rect(s3, cx, cy, card_w, card_h, WHITE, line=DIVIDER)
    # accent top bar
    add_rect(s3, cx, cy, card_w, Inches(0.12), accent)
    # Icon on left
    add_image(s3, ICONS / icon_f, cx + Inches(0.3), cy + Inches(0.4),
              w=Inches(1.35), h=Inches(1.35))
    # Number
    add_text(s3, cx + Inches(1.85), cy + Inches(0.4), Inches(1.0), Inches(0.35),
             num, size=14, bold=True, color=accent)
    # Title
    add_text(s3, cx + Inches(1.85), cy + Inches(0.75), card_w - Inches(2.1), Inches(0.5),
             title, size=20, bold=True, color=NAVY)
    # Body
    add_text(s3, cx + Inches(1.85), cy + Inches(1.25), card_w - Inches(2.1), card_h - Inches(1.4),
             body, size=11.5, color=TEXT_MID, line_spacing=1.4)

page_footer(s3, "ESP32 IoT Enhancement · Solution Overview", "Slide 3 of 11")

add_notes(s3, """'We solved both problems with four additions.

One — SIM Module. The device can now use a mobile cellular network directly. It no longer depends only on WiFi.

Two — SD Card Backup. Whenever the network goes down, the SD card activates automatically and stores every data ping locally. No more permanent data loss.

Three — RTC Module, which is a real-time clock. It keeps accurate date and time continuously, so every ping stored on the SD card gets a correct timestamp.

Four — a Web Server. We built a browser-based configuration portal. Users can now configure the device from any web browser. No Arduino editing, no firmware reflash, no developer required.

The next four slides go into each one.'

Timing: 60 seconds.""")


# ============================================================
#  SLIDE 4 — SIM MODULE
# ============================================================
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)
page_header(s4, "Feature 01  ·  SIM Module", "Cellular connectivity for unstable WiFi locations")

# Left — large icon with soft badge background
icon_bg_x = Inches(0.9); icon_bg_y = Inches(2.0)
icon_bg_size = Inches(4.0)
add_rounded(s4, icon_bg_x, icon_bg_y, icon_bg_size, icon_bg_size,
            RGBColor(0xEB, 0xF8, 0xF7), radius=0.18)
add_image(s4, ICONS / "wifi.png",
          icon_bg_x + Inches(0.5), icon_bg_y + Inches(0.5),
          w=Inches(3.0), h=Inches(3.0))

# caption under icon
add_text(s4, icon_bg_x, icon_bg_y + icon_bg_size + Inches(0.25),
         icon_bg_size, Inches(0.4),
         "Mobile SIM Network", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s4, icon_bg_x, icon_bg_y + icon_bg_size + Inches(0.65),
         icon_bg_size, Inches(0.3),
         "Cellular (A7670C 4G LTE)", size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)

# Right — benefits panel
right_x = Inches(5.4)
right_w = Inches(7.4)
add_text(s4, right_x, Inches(1.8), right_w, Inches(0.4),
         "BENEFITS", size=11, bold=True, color=TEAL_DARK)
add_text(s4, right_x, Inches(2.15), right_w, Inches(0.6),
         "What the SIM Module delivers.", size=22, bold=True, color=NAVY, line_spacing=1.1)

sim_benefits = [
    ("Mobile SIM network, directly.",
     "Device can now connect through a cellular SIM — no longer dependent on WiFi."),
    ("Reduced WiFi dependency.",
     "Even if local WiFi fails, the device stays connected over the mobile network."),
    ("Works in remote or unstable areas.",
     "Sites with weak or no WiFi now have a reliable connectivity path."),
    ("Fewer transmission failures.",
     "Telemetry reaches the cloud even when the primary network is degraded."),
    ("Higher uptime and reliability.",
     "Overall device availability improves across all deployment conditions."),
]
cur_y = Inches(2.95)
for title, body in sim_benefits:
    # dot
    add_rect(s4, right_x, cur_y + Inches(0.1), Inches(0.12), Inches(0.12), TEAL)
    add_text(s4, right_x + Inches(0.3), cur_y, right_w - Inches(0.3), Inches(0.32),
             title, size=13, bold=True, color=NAVY)
    add_text(s4, right_x + Inches(0.3), cur_y + Inches(0.32), right_w - Inches(0.3), Inches(0.4),
             body, size=11, color=TEXT_MID, line_spacing=1.35)
    cur_y += Inches(0.75)

page_footer(s4, "ESP32 IoT Enhancement · SIM Module", "Slide 4 of 11")

add_notes(s4, """'First addition — the SIM Module.

The device can now connect to the internet through a mobile cellular SIM, directly. This removes the single-point-of-failure we had with WiFi-only connectivity.

Five concrete benefits.

One — it uses the mobile network directly, just like a smartphone.

Two — it reduces dependency on WiFi. Even if local WiFi fails, the device stays connected.

Three — it works in remote or unstable areas where WiFi either doesn't exist or isn't reliable.

Four — it minimizes transmission failures. Our telemetry reaches the cloud under a wider range of field conditions.

Five — overall uptime and reliability goes up. Significantly.'

Timing: 70 seconds.""")


# ============================================================
#  SLIDE 5 — SD CARD BACKUP
# ============================================================
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)
page_header(s5, "Feature 02  ·  SD Card Backup", "Automatic offline storage during network outages")

# Left — icon with badge
icon_bg_x = Inches(0.9); icon_bg_y = Inches(2.0)
icon_bg_size = Inches(4.0)
add_rounded(s5, icon_bg_x, icon_bg_y, icon_bg_size, icon_bg_size,
            RGBColor(0xEB, 0xF8, 0xF7), radius=0.18)
add_image(s5, ICONS / "sdcard.png",
          icon_bg_x + Inches(0.5), icon_bg_y + Inches(0.5),
          w=Inches(3.0), h=Inches(3.0))

add_text(s5, icon_bg_x, icon_bg_y + icon_bg_size + Inches(0.25),
         icon_bg_size, Inches(0.4),
         "Local Offline Storage", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s5, icon_bg_x, icon_bg_y + icon_bg_size + Inches(0.65),
         icon_bg_size, Inches(0.3),
         "FAT32  ·  Chronological replay", size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)

# Right — benefits + mini flow
right_x = Inches(5.4)
right_w = Inches(7.4)
add_text(s5, right_x, Inches(1.8), right_w, Inches(0.4),
         "HOW IT WORKS", size=11, bold=True, color=TEAL_DARK)
add_text(s5, right_x, Inches(2.15), right_w, Inches(0.6),
         "Zero data loss during outages.", size=22, bold=True, color=NAVY, line_spacing=1.1)

sd_benefits = [
    ("Activates automatically.",
     "When the internet or network goes down, SD backup kicks in without any manual step."),
    ("Stores every data ping locally.",
     "Every reading is written to the SD card with its timestamp, safely preserved."),
    ("Prevents data loss during outages.",
     "No matter how long the network is down, nothing is dropped."),
    ("Chronological replay on recovery.",
     "When the network returns, cached data is sent to the cloud in the correct order."),
]
cur_y = Inches(2.95)
for title, body in sd_benefits:
    add_rect(s5, right_x, cur_y + Inches(0.1), Inches(0.12), Inches(0.12), TEAL)
    add_text(s5, right_x + Inches(0.3), cur_y, right_w - Inches(0.3), Inches(0.32),
             title, size=13, bold=True, color=NAVY)
    add_text(s5, right_x + Inches(0.3), cur_y + Inches(0.32), right_w - Inches(0.3), Inches(0.4),
             body, size=11, color=TEXT_MID, line_spacing=1.35)
    cur_y += Inches(0.75)

# Mini flow diagram at bottom
flow_y = Inches(6.35)
flow_x = Inches(0.6)
flow_w = Inches(12.15)
add_rounded(s5, flow_x, flow_y, flow_w, Inches(0.55), NAVY, radius=0.3)
add_text(s5, flow_x + Inches(0.25), flow_y + Inches(0.05), flow_w - Inches(0.5), Inches(0.25),
         "OFFLINE FLOW", size=9, bold=True, color=TEAL)
add_text(s5, flow_x + Inches(0.25), flow_y + Inches(0.25), flow_w - Inches(0.5), Inches(0.3),
         "Network drop   →   SD backup activates   →   Every ping stored locally   →   Network back   →   Chronological replay to cloud",
         size=11, bold=True, color=WHITE)

page_footer(s5, "ESP32 IoT Enhancement · SD Card Backup", "Slide 5 of 11")

add_notes(s5, """'Second addition — SD Card Backup.

Whenever the internet or network connection goes down, the SD card backup activates automatically. There's no manual step, no user intervention — the device detects the outage and switches to local storage.

Every data ping gets stored locally on the SD card with its timestamp.

When the network comes back, the device replays all the cached data to the cloud in chronological order — so the server sees the readings in the exact sequence they were captured.

The bottom bar shows the flow — network drop, SD activates, store locally, network recovers, replay. No data is lost, no matter how long the outage lasts.'

Timing: 70 seconds.""")


# ============================================================
#  SLIDE 6 — RTC MODULE
# ============================================================
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)
page_header(s6, "Feature 03  ·  RTC Module", "Real-Time Clock for accurate timestamps")

# Left — icon
icon_bg_x = Inches(0.9); icon_bg_y = Inches(2.0)
icon_bg_size = Inches(4.0)
add_rounded(s6, icon_bg_x, icon_bg_y, icon_bg_size, icon_bg_size,
            RGBColor(0xFD, 0xF5, 0xEC), radius=0.18)
add_image(s6, ICONS / "refresh.png",
          icon_bg_x + Inches(0.5), icon_bg_y + Inches(0.5),
          w=Inches(3.0), h=Inches(3.0))

add_text(s6, icon_bg_x, icon_bg_y + icon_bg_size + Inches(0.25),
         icon_bg_size, Inches(0.4),
         "DS3231 RTC", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s6, icon_bg_x, icon_bg_y + icon_bg_size + Inches(0.65),
         icon_bg_size, Inches(0.3),
         "Accurate time, even offline", size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)

# Right — purpose
right_x = Inches(5.4)
right_w = Inches(7.4)
add_text(s6, right_x, Inches(1.8), right_w, Inches(0.4),
         "PURPOSE", size=11, bold=True, color=AMBER_DARK)
add_text(s6, right_x, Inches(2.15), right_w, Inches(0.6),
         "Every ping gets the right time.", size=22, bold=True, color=NAVY, line_spacing=1.1)

rtc_points = [
    ("RTC = Real-Time Clock Module.",
     "Dedicated hardware that tracks the current date and time continuously."),
    ("Accurate time, even offline.",
     "The RTC keeps correct time even when the device is disconnected from the network."),
    ("Every stored ping gets an exact timestamp.",
     "When SD backup is active, each reading is tagged with the precise time of capture."),
    ("Correct historical records.",
     "Server-side logs, reports, and compliance data reflect actual capture time — not arrival time."),
    ("Traceability across outages.",
     "You can always tell when each reading was taken, regardless of when it reached the cloud."),
]
cur_y = Inches(2.95)
for title, body in rtc_points:
    add_rect(s6, right_x, cur_y + Inches(0.1), Inches(0.12), Inches(0.12), AMBER)
    add_text(s6, right_x + Inches(0.3), cur_y, right_w - Inches(0.3), Inches(0.32),
             title, size=13, bold=True, color=NAVY)
    add_text(s6, right_x + Inches(0.3), cur_y + Inches(0.32), right_w - Inches(0.3), Inches(0.4),
             body, size=11, color=TEXT_MID, line_spacing=1.35)
    cur_y += Inches(0.72)

page_footer(s6, "ESP32 IoT Enhancement · RTC Module", "Slide 6 of 11")

add_notes(s6, """'Third addition — the RTC Module.

RTC stands for Real-Time Clock. It's a dedicated piece of hardware that tracks the current date and time continuously — independent of the network and independent of whether the device has power to the main processor.

Why does this matter? Because when the SD card is buffering data during a network outage, we need to know exactly when each reading was captured — not when it eventually reaches the cloud.

The RTC makes sure every ping stored on the SD card gets an exact, accurate timestamp. This means our historical records, compliance reports, and analytics all reflect the real time of capture.

It also means we have full traceability across any outage — we can always tell exactly when each reading happened.'

Timing: 65 seconds.""")


# ============================================================
#  SLIDE 7 — WEB SERVER
# ============================================================
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)
page_header(s7, "Feature 04  ·  Web Server", "Browser-based configuration for easy deployment")

# Left — icon
icon_bg_x = Inches(0.9); icon_bg_y = Inches(2.0)
icon_bg_size = Inches(4.0)
add_rounded(s7, icon_bg_x, icon_bg_y, icon_bg_size, icon_bg_size,
            RGBColor(0xFD, 0xF5, 0xEC), radius=0.18)
add_image(s7, ICONS / "globe.png",
          icon_bg_x + Inches(0.5), icon_bg_y + Inches(0.5),
          w=Inches(3.0), h=Inches(3.0))

add_text(s7, icon_bg_x, icon_bg_y + icon_bg_size + Inches(0.25),
         icon_bg_size, Inches(0.4),
         "Web Configuration Portal", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s7, icon_bg_x, icon_bg_y + icon_bg_size + Inches(0.65),
         icon_bg_size, Inches(0.3),
         "Browser-based  ·  No code, no reflash", size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)

# Right — BEFORE / AFTER comparison
right_x = Inches(5.4)
right_w = Inches(7.4)
add_text(s7, right_x, Inches(1.8), right_w, Inches(0.4),
         "BEFORE  ·  AFTER", size=11, bold=True, color=AMBER_DARK)
add_text(s7, right_x, Inches(2.15), right_w, Inches(0.6),
         "How configuration changed.", size=22, bold=True, color=NAVY, line_spacing=1.1)

# Before column
col_w = Inches(3.6)
before_x = right_x
after_x = right_x + col_w + Inches(0.2)
col_y = Inches(3.05)
col_h = Inches(3.7)

# BEFORE
add_rect(s7, before_x, col_y, col_w, col_h, WHITE, line=DIVIDER)
add_rect(s7, before_x, col_y, col_w, Inches(0.4), ALERT)
add_text(s7, before_x + Inches(0.25), col_y + Inches(0.07), col_w - Inches(0.5), Inches(0.3),
         "BEFORE", size=11, bold=True, color=WHITE)
before_items = [
    "Open Arduino IDE",
    "Edit source code manually",
    "Recompile firmware",
    "Reflash the ESP32",
    "Requires a developer",
    "Slow, technical, error-prone",
]
by = col_y + Inches(0.6)
for it in before_items:
    add_text(s7, before_x + Inches(0.25), by, col_w - Inches(0.5), Inches(0.3),
             "×  " + it, size=11.5, color=TEXT_MID, line_spacing=1.3)
    by += Inches(0.48)

# AFTER
add_rect(s7, after_x, col_y, col_w, col_h, WHITE, line=DIVIDER)
add_rect(s7, after_x, col_y, col_w, Inches(0.4), SUCCESS)
add_text(s7, after_x + Inches(0.25), col_y + Inches(0.07), col_w - Inches(0.5), Inches(0.3),
         "AFTER", size=11, bold=True, color=WHITE)
after_items = [
    "Open browser",
    "Navigate to device portal",
    "Change settings in UI",
    "Click save",
    "Anyone can do it",
    "Fast, user-friendly",
]
ay = col_y + Inches(0.6)
for it in after_items:
    add_text(s7, after_x + Inches(0.25), ay, col_w - Inches(0.5), Inches(0.3),
             "✓  " + it, size=11.5, color=TEXT_DARK, bold=True, line_spacing=1.3)
    ay += Inches(0.48)

page_footer(s7, "ESP32 IoT Enhancement · Web Server", "Slide 7 of 11")

add_notes(s7, """'Fourth addition — the Web Server.

Before this, changing any configuration on the device meant opening the Arduino IDE, editing the source code, recompiling, and reflashing the firmware. You needed a developer, a laptop with the build environment, and time.

We built a browser-based configuration portal. Now anyone — technical or not — can open a browser, connect to the device, change settings in a clean UI, and click save. That's it.

On the right you can see the before-and-after. Before took six technical steps. After takes four simple clicks, and anyone can do it.

This is the change that makes the device genuinely deployable at scale — because rolling out a hundred sites no longer bottlenecks on developer availability.'

Timing: 70 seconds.""")


# ============================================================
#  SLIDE 8 — WEB PORTAL (LIVE VIEW WITH SCREENSHOTS)
# ============================================================
s_web = prs.slides.add_slide(BLANK)
add_rect(s_web, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)
page_header(s_web, "Feature 04  ·  Web Portal Live View",
            "Modern, browser-based configuration interface")

# Screenshot layout — two images side by side
shot_w = Inches(5.95)
shot_h = Inches(3.15)  # preserves aspect ratio of the source screenshots
shot_y = Inches(1.75)
shot_gap = Inches(0.25)
left_shot_x = Inches(0.6)
right_shot_x = left_shot_x + shot_w + shot_gap

# LEFT — Overview screenshot with frame
# Shadow behind
add_rect(s_web, left_shot_x + Emu(50000), shot_y + Emu(50000),
         shot_w, shot_h, RGBColor(0xCF, 0xD5, 0xDC))
# White frame/padding
add_rounded(s_web, left_shot_x, shot_y, shot_w, shot_h, WHITE,
            line=DIVIDER, radius=0.02)
# Actual screenshot (with some padding inside the frame)
pad = Inches(0.08)
add_image(s_web, ASSETS / "screenshots" / "webserver_overview.png",
          left_shot_x + pad, shot_y + pad,
          w=shot_w - pad*2, h=shot_h - pad*2)

# LEFT caption block
lcap_y = shot_y + shot_h + Inches(0.15)
add_rect(s_web, left_shot_x, lcap_y, Inches(0.15), Inches(0.9), TEAL)
add_text(s_web, left_shot_x + Inches(0.3), lcap_y, shot_w - Inches(0.4), Inches(0.4),
         "01  ·  Dashboard", size=15, bold=True, color=NAVY)
add_text(s_web, left_shot_x + Inches(0.3), lcap_y + Inches(0.4),
         shot_w - Inches(0.4), Inches(0.5),
         "Real-time system status — firmware version, MAC address, uptime, heap usage, active tasks.",
         size=11, color=TEXT_MID, line_spacing=1.35)

# RIGHT — Sensors screenshot with frame
add_rect(s_web, right_shot_x + Emu(50000), shot_y + Emu(50000),
         shot_w, shot_h, RGBColor(0xCF, 0xD5, 0xDC))
add_rounded(s_web, right_shot_x, shot_y, shot_w, shot_h, WHITE,
            line=DIVIDER, radius=0.02)
add_image(s_web, ASSETS / "screenshots" / "webserver_sensors.png",
          right_shot_x + pad, shot_y + pad,
          w=shot_w - pad*2, h=shot_h - pad*2)

# RIGHT caption block
rcap_y = shot_y + shot_h + Inches(0.15)
add_rect(s_web, right_shot_x, rcap_y, Inches(0.15), Inches(0.9), AMBER)
add_text(s_web, right_shot_x + Inches(0.3), rcap_y, shot_w - Inches(0.4), Inches(0.4),
         "02  ·  Sensor Management", size=15, bold=True, color=NAVY)
add_text(s_web, right_shot_x + Inches(0.3), rcap_y + Inches(0.4),
         shot_w - Inches(0.4), Inches(0.5),
         "Add, edit, test and delete Modbus sensors from the browser — no code, no reflash.",
         size=11, color=TEXT_MID, line_spacing=1.35)

# Bottom navigation pill strip — showing all 6 portal sections
nav_y = Inches(6.5)
add_rounded(s_web, Inches(0.6), nav_y, Inches(12.15), Inches(0.5), NAVY, radius=0.3)
add_text(s_web, Inches(0.9), nav_y + Inches(0.05), Inches(11.6), Inches(0.2),
         "PORTAL SECTIONS", size=8, bold=True, color=TEAL)
add_text(s_web, Inches(0.9), nav_y + Inches(0.22), Inches(11.6), Inches(0.3),
         "Overview  ·  Network Config  ·  Azure IoT Hub  ·  Modbus Sensors  ·  Write Operations  ·  OTA Update",
         size=12, bold=True, color=WHITE)

page_footer(s_web, "ESP32 IoT Enhancement · Web Portal", "Slide 8 of 11")

add_notes(s_web, """Show the actual product. Let the screenshots speak.

'This is what the portal actually looks like — running live on the device.

On the left — the Dashboard. Real-time system status. Firmware version, MAC address, uptime, heap usage, active tasks. Operations teams can check device health at a glance without needing a terminal or a developer.

On the right — Sensor Management. We can add a new Modbus sensor, configure its unit ID, slave ID, register, data type, byte order, scale factor — all from the browser. Then test the connection with a single click. Edit or delete at any time.

The bar at the bottom shows the six portal sections — Overview, Network Config, Azure IoT Hub, Modbus Sensors, Write Operations, OTA Update. Every device setting is reachable from here.

This portal is what replaced the Arduino-IDE-and-reflash workflow.'

Timing: 75 seconds. Let the room look at the screenshots — don't rush.""")


# ============================================================
#  SLIDE 9 — OVERALL BUSINESS IMPACT
# ============================================================
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)
page_header(s8, "Overall Business Impact", "How the transformation helps the business")

add_text(s8, Inches(0.6), Inches(1.55), Inches(12), Inches(0.4),
         "Six outcomes — direct results of the four additions above.",
         size=13, color=TEXT_MID)

# 6 impact tiles — 3x2 grid
tile_x0 = Inches(0.6)
tile_y0 = Inches(2.15)
tile_w = Inches(4.0)
tile_h = Inches(2.3)
tile_gx = Inches(0.1)
tile_gy = Inches(0.15)

impacts_biz = [
    ("Major reduction\nin data loss.",
     "SD backup + SIM failover keep data flowing through outages.",
     TEAL, "sdcard.png"),
    ("Improved connectivity\nreliability.",
     "Dual network removes the single-point-of-failure from WiFi.",
     TEAL, "wifi.png"),
    ("Easier field deployment.",
     "Plug-and-play commissioning via the browser — no laptops, no code.",
     TEAL, "globe.png"),
    ("Reduced technical\ndependency.",
     "Non-technical users can configure, monitor, and adjust settings.",
     AMBER, "p_laptop.png"),
    ("Faster customer\nsupport.",
     "Configuration changes happen remotely, without site visits.",
     AMBER, "refresh.png"),
    ("Better user experience.",
     "Simple UI and reliable behavior across all network conditions.",
     AMBER, "p_eye.png"),
]

for i, (title, body, accent, icon_f) in enumerate(impacts_biz):
    col = i % 3; row = i // 3
    tx = tile_x0 + col * (tile_w + tile_gx)
    ty = tile_y0 + row * (tile_h + tile_gy)
    # shadow
    add_rect(s8, tx + Emu(40000), ty + Emu(40000), tile_w, tile_h, RGBColor(0xDD, 0xE1, 0xE7))
    # body
    add_rect(s8, tx, ty, tile_w, tile_h, WHITE, line=DIVIDER)
    # accent top bar
    add_rect(s8, tx, ty, tile_w, Inches(0.1), accent)
    # icon
    add_image(s8, ICONS / icon_f, tx + Inches(0.3), ty + Inches(0.35),
              w=Inches(0.75), h=Inches(0.75))
    # title
    add_text(s8, tx + Inches(1.2), ty + Inches(0.35), tile_w - Inches(1.4), Inches(0.85),
             title, size=14, bold=True, color=NAVY, line_spacing=1.15)
    # body
    add_text(s8, tx + Inches(0.3), ty + Inches(1.3), tile_w - Inches(0.5), tile_h - Inches(1.4),
             body, size=10.5, color=TEXT_MID, line_spacing=1.4)

page_footer(s8, "ESP32 IoT Enhancement · Business Impact", "Slide 9 of 11")

add_notes(s8, """'Adding those four features gave us six direct business outcomes.

One — major reduction in data loss. The SD backup and SIM failover together mean data keeps flowing even when networks fail.

Two — improved connectivity reliability. Dual network takes WiFi out of being a single point of failure.

Three — easier field deployment. Commissioning a new site is now a browser-and-settings operation, not a firmware operation.

Four — reduced technical dependency. Non-technical users can configure, monitor, and adjust the device.

Five — faster customer support. Configuration changes happen remotely, no site visits needed.

Six — better user experience across the board. Simpler UI, more reliable behavior.'

Timing: 75 seconds.""")


# ============================================================
#  SLIDE 10 — IN-HOUSE vs EXTERNAL (AQUAGEN vs TELTONIKA)
# ============================================================
s_compare = prs.slides.add_slide(BLANK)
add_rect(s_compare, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_LIGHT)
page_header(s_compare, "In-House vs External",
            "Aquagen vs Teltonika — feature-by-feature comparison")

# Subtitle line
add_text(s_compare, Inches(0.6), Inches(1.55), Inches(12.15), Inches(0.4),
         "Feature parity on core capabilities — with ownership, direct Azure integration, and no vendor subscription.",
         size=12, color=TEXT_MID, line_spacing=1.3)

# Two product cards side by side
card_w = Inches(6.05)
card_h = Inches(4.95)
card_y = Inches(2.1)
card_gap = Inches(0.15)
left_card_x = Inches(0.6)
right_card_x = left_card_x + card_w + card_gap


def product_card(slide, x, y, w, h, *, title_text, subtitle_text, tag_text,
                 accent_color, row_items):
    """Draw a product comparison card."""
    # Shadow
    add_rect(slide, x + Emu(45000), y + Emu(45000), w, h, RGBColor(0xDD, 0xE1, 0xE7))
    # Body
    add_rect(slide, x, y, w, h, WHITE, line=DIVIDER)
    # Header band
    add_rect(slide, x, y, w, Inches(1.0), accent_color)
    # Product name
    add_text(slide, x + Inches(0.4), y + Inches(0.15), w - Inches(0.8), Inches(0.35),
             tag_text.upper(), size=10, bold=True, color=WHITE)
    add_text(slide, x + Inches(0.4), y + Inches(0.45), w - Inches(0.8), Inches(0.5),
             title_text, size=22, bold=True, color=WHITE)
    add_text(slide, x + Inches(0.4), y + Inches(0.78), w - Inches(0.8), Inches(0.3),
             subtitle_text, size=10, color=RGBColor(0xE4, 0xF8, 0xF7))

    # Feature rows
    row_y = y + Inches(1.2)
    row_h = Inches(0.38)
    for i, (label, value) in enumerate(row_items):
        # alternating shading
        if i % 2 == 0:
            add_rect(slide, x + Inches(0.2), row_y, w - Inches(0.4), row_h,
                     RGBColor(0xFA, 0xFB, 0xFC))
        # label
        add_text(slide, x + Inches(0.35), row_y + Inches(0.09),
                 Inches(2.1), Inches(0.25),
                 label, size=9.5, bold=True, color=TEXT_MID)
        # value
        add_text(slide, x + Inches(2.55), row_y + Inches(0.07),
                 w - Inches(2.75), Inches(0.3),
                 value, size=10.5, color=NAVY, line_spacing=1.25)
        row_y += row_h


# AQUAGEN (in-house) — teal accent
aquagen_rows = [
    ("Platform",          "ESP32-WROOM-32 · Custom hardware"),
    ("Connectivity",      "WiFi + 4G LTE (A7670C cellular)"),
    ("Industrial bus",    "RS485 Modbus RTU · 10 sensors per unit"),
    ("Offline storage",   "SD Card · automatic + chronological replay"),
    ("Timestamps",        "DS3231 RTC · accurate even offline"),
    ("Cloud integration", "Azure IoT Hub · MQTT/TLS direct"),
    ("Remote management", "Azure Device Twin · built-in · no subscription"),
    ("Configuration",     "Custom web portal · branded UI"),
    ("OTA updates",       "Signed firmware · dual-slot rollback"),
    ("Customization",     "Full source control · any change possible"),
    ("Vendor dependency", "None · built in-house at FluxGen"),
]
product_card(s_compare, left_card_x, card_y, card_w, card_h,
             title_text="Aquagen", subtitle_text="FluxGen in-house IoT gateway",
             tag_text="IN-HOUSE · BUILT BY US",
             accent_color=TEAL_DARK, row_items=aquagen_rows)

# TELTONIKA (external) — navy/amber accent
teltonika_rows = [
    ("Platform",          "RUT956 / TRB246 · vendor-built router"),
    ("Connectivity",      "WiFi + 4G LTE · Dual-SIM failover"),
    ("Industrial bus",    "RS232 + RS485 · digital I/O + relay"),
    ("Offline storage",   "microSD (requires configuration)"),
    ("Timestamps",        "System clock · NTP-synced"),
    ("Cloud integration", "Via MQTT broker or Teltonika RMS"),
    ("Remote management", "Teltonika RMS · paid subscription"),
    ("Configuration",     "RutOS web UI · vendor interface"),
    ("OTA updates",       "Vendor-supplied firmware updates"),
    ("Customization",     "Limited to RutOS parameters"),
    ("Vendor dependency", "Teltonika (Lithuania-based vendor)"),
]
product_card(s_compare, right_card_x, card_y, card_w, card_h,
             title_text="Teltonika", subtitle_text="External procurement · RUT956 / TRB246 series",
             tag_text="EXTERNAL · WE BUY",
             accent_color=AMBER_DARK, row_items=teltonika_rows)

# Bottom takeaway pill
takeaway_y = Inches(7.15)
add_rounded(s_compare, Inches(0.6), takeaway_y - Inches(0.03), Inches(12.15), Inches(0.4),
            NAVY, radius=0.4)
add_text(s_compare, Inches(0.9), takeaway_y, Inches(11.5), Inches(0.35),
         "We match them on core features — and own the stack, the roadmap, and the cloud integration.",
         size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

page_footer(s_compare, "ESP32 IoT Enhancement · In-House vs External",
            "Slide 10 of 11")

add_notes(s_compare, """'Let me put this in context. Teltonika is what we currently buy externally for industrial IoT deployments — specifically their RUT956 and TRB246 series. Good products. Widely used.

On the left — Aquagen, our in-house build. On the right — Teltonika, what we buy.

Look at the feature rows. On connectivity — both do WiFi plus 4G LTE. On industrial bus — both do RS485 Modbus. On offline storage — both support SD-card caching, though on Teltonika it requires specific configuration. On cloud — Teltonika integrates via their RMS subscription service or an MQTT broker; Aquagen connects directly to Azure IoT Hub.

The real differences are at the bottom. Customization — Aquagen is fully in-house, so we can change anything. Teltonika is locked to what RutOS supports. Vendor dependency — Aquagen has none, Teltonika ties us to a Lithuania-based vendor and their subscription model.

The bottom line — we match them on core features, but we own the stack, the roadmap, and the cloud integration.'

Timing: 70 seconds. This slide is where the business case for 'in-house' lands.""")


# ============================================================
#  SLIDE 11 — FINAL SUMMARY & CLOSING
# ============================================================
s9 = prs.slides.add_slide(BLANK)
add_image(s9, ASSETS / "bg_closing.png", Inches(0), Inches(0), SLIDE_W, SLIDE_H)

add_text(s9, Inches(1.5), Inches(0.8), Inches(10.3), Inches(0.4),
         "FINAL SUMMARY", size=12, bold=True, color=TEAL)

# Main closing statement
add_text(s9, Inches(1.5), Inches(1.5), Inches(10.3), Inches(2.0),
         "From Data Loss\nto Smart Reliability.",
         size=48, bold=True, color=WHITE, line_spacing=1.1)

# Sub-statement
add_text(s9, Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.8),
         "A Complete IoT Upgrade.",
         size=28, color=TEAL, line_spacing=1.2)

# Divider
add_rect(s9, Inches(1.5), Inches(4.75), Inches(1.2), Inches(0.04), AMBER)

# What we transformed
add_text(s9, Inches(1.5), Inches(4.95), Inches(10.3), Inches(0.35),
         "The device transformed from a WiFi-dependent system into a smart hybrid IoT solution.",
         size=14, color=WHITE, bold=True)

add_text(s9, Inches(1.5), Inches(5.4), Inches(10.3), Inches(0.35),
         "Dual network capability  ·  Offline data backup  ·  Accurate timestamp logging  ·  Easy web-based control",
         size=12, color=TEXT_LIGHT)

# Thank you block
add_text(s9, Inches(1.5), Inches(6.55), Inches(10.3), Inches(0.35),
         "Thank you",
         size=16, bold=True, color=WHITE)
add_text(s9, Inches(1.5), Inches(6.95), Inches(10.3), Inches(0.3),
         "Y Chaithanya Reddy  ·  FluxGen  ·  IoT Platform",
         size=10, color=TEXT_LIGHT)

add_notes(s9, """Don't read the slide. Deliver from memory, quietly, with weight.

'To summarize —

This project took a WiFi-dependent ESP32 device and transformed it into a smart hybrid IoT solution.

Four additions made the difference — SIM module for cellular connectivity, SD card for offline backup, RTC for accurate timestamps, and a web server for easy configuration.

The result — major reduction in data loss, improved reliability, easier deployment, and a much better user experience.

From data loss to smart reliability. A complete IoT upgrade.

Thank you.'

Then stop talking. Pause. Let silence work. Five seconds minimum.""")


# ============================================================
#  TRANSITIONS (fade on every slide)
# ============================================================
print("Injecting transitions...")
for s in prs.slides:
    add_slide_transition(s, transition_type="fade")


# ============================================================
#  SAVE
# ============================================================
out_path = r"C:\Users\chath\Documents\Python code\modbus_iot_gateway\presentation\FluxGen_IoT_Gateway_v2.0_Pitch.pptx"
prs.save(out_path)
print(f"Generated: {out_path}")
print(f"Slides: {len(prs.slides)}")
